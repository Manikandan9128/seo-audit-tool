"""Builds a client-facing PPTX audit report styled after the Cyces deck format:
dark top bar, blue section-title band, light-gray body, white content cards."""

from __future__ import annotations

import re
import threading
from collections import Counter
from io import BytesIO
from urllib.parse import urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

from app.services.keyword_relevance_service import (
    _KEYWORD_PAGE_CATEGORIES,
    _brand_token,
    _classify_keyword_page_category,
    _is_branded_keyword,
)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DARK = RGBColor(0x14, 0x17, 0x1C)
# Cyces' own brand red (confirmed from cyces.co's compiled CSS, #FF0000) —
# the report's default accent whenever a client's site has no extractable
# brand color of its own. A client's real brand color (via brand_color_hex)
# always overrides this.
DEFAULT_ACCENT = RGBColor(0xFF, 0x00, 0x00)
ORANGE = RGBColor(0xF2, 0x8C, 0x28)
BODY_BG = RGBColor(0xF4, 0xF5, 0xF7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x14, 0x17, 0x1C)
TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)
GOOD = RGBColor(0x1F, 0x9D, 0x66)
WARN = RGBColor(0xE0, 0x8E, 0x1D)
BAD = RGBColor(0xD1, 0x3B, 0x3B)
ROW_ALT = RGBColor(0xF8, 0xF9, 0xFA)
CARD_BORDER = RGBColor(0xE6, 0xE8, 0xEB)
HEADER_ROW_BG = RGBColor(0xF1, 0xF3, 0xF5)
HEADER_ROW_TEXT = RGBColor(0x4A, 0x4A, 0x4A)
CHIP_BG = RGBColor(0xEA, 0xF4, 0xF8)

# Mutable so build_report can theme a single generation run with the client's
# own brand color (scraped from their site) without threading a param through
# every slide-drawing function. Thread-local (not a plain dict) because
# report requests run concurrently in Starlette's thread pool — a shared
# global here let one in-flight build's accent color/footer bleed into
# another's slides mid-build, producing a deck with mismatched colors
# per slide whenever two reports were generated around the same time.
class _ThemeStore(threading.local):
    def __init__(self):
        self.accent = DEFAULT_ACCENT
        self.footer = ""


class _ThemeProxy:
    def __getitem__(self, key):
        return getattr(_theme_store, key)

    def __setitem__(self, key, value):
        setattr(_theme_store, key, value)


_theme_store = _ThemeStore()
_theme = _ThemeProxy()


def _accent() -> RGBColor:
    return _theme["accent"]


def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _textbox(slide, left, top, width, height, text, size=14, bold=False, color=TEXT_DARK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _body_bg(slide):
    rect = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    _fill(rect, BODY_BG)
    rect.shadow.inherit = False
    return rect


def _content_header(slide, title, eyebrow=None):
    """Thin dark top strip + optional eyebrow tag + bold section title +
    full-width hairline accent rule, used on every content slide. Compact by
    design so slides carry more content below the fold. Also draws the footer."""
    _body_bg(slide)
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.1))
    _fill(bar, DARK)
    bar.shadow.inherit = False

    title_y = Inches(0.22)
    if eyebrow:
        _textbox(slide, Inches(0.4), Inches(0.18), Inches(8), Inches(0.28), eyebrow.upper(), size=10.5, bold=True, color=_accent())
        title_y = Inches(0.44)
    _textbox(slide, Inches(0.4), title_y, Inches(10.5), Inches(0.5), title, size=23, bold=True, color=TEXT_DARK)

    rule = slide.shapes.add_shape(1, Inches(0.4), Inches(0.92), SLIDE_W - Inches(0.8), Pt(1))
    _fill(rule, CARD_BORDER)
    rule.shadow.inherit = False
    accent_tick = slide.shapes.add_shape(1, Inches(0.4), Inches(0.9), Inches(0.55), Pt(3))
    _fill(accent_tick, _accent())
    accent_tick.shadow.inherit = False
    _footer(slide)


def _footer(slide):
    if not _theme["footer"]:
        return
    _textbox(
        slide, Inches(0.4), SLIDE_H - Inches(0.4), Inches(8), Inches(0.3),
        _theme["footer"], size=9, color=TEXT_MUTED,
    )


def _card(slide, left, top, width, height):
    """White rounded panel with a hairline border and a soft drop shadow,
    matching the Semrush/PSI-style panels in the sample deck."""
    card = slide.shapes.add_shape(5, left, top, width, height)  # rounded rectangle
    try:
        card.adjustments[0] = 0.045
    except (IndexError, AttributeError):
        pass
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = CARD_BORDER
    card.line.width = Pt(0.75)
    card.shadow.inherit = False
    _soft_shadow(card)
    return card


def _soft_shadow(shape):
    """Subtle drop shadow via the raw XML shadow effect — python-pptx has no
    high-level API for outer shadows.

    `shape.shadow.inherit = False` (set by callers before this) already
    inserts its own empty <a:effectLst/> to suppress the theme's inherited
    shadow. Appending a second one on top of it is invalid per the OOXML
    schema (effectLst may appear at most once) — PowerPoint silently
    "repaired" every file built this way by stripping the invalid element,
    taking that shape's fill/shadow styling down with it. Removing any
    existing effectLst first keeps this to the one populated element the
    schema allows."""
    from pptx.oxml.ns import qn

    sp_pr = shape._element.spPr
    for existing in sp_pr.findall(qn("a:effectLst")):
        sp_pr.remove(existing)
    effect_lst = sp_pr.makeelement(qn("a:effectLst"), {})
    outer_shdw = sp_pr.makeelement(
        qn("a:outerShdw"),
        {"blurRad": "90000", "dist": "20000", "dir": "5400000", "rotWithShape": "0"},
    )
    clr = sp_pr.makeelement(qn("a:srgbClr"), {"val": "1A1A1A"})
    alpha = sp_pr.makeelement(qn("a:alpha"), {"val": "12000"})
    clr.append(alpha)
    outer_shdw.append(clr)
    effect_lst.append(outer_shdw)
    sp_pr.append(effect_lst)


def _chip_row(slide, left, top, items, max_width, size=11, bg=None, fg=None):
    """Compact wrapping row of pill-shaped tags — denser and more modern than
    a comma-joined text line, and self-wraps onto multiple rows if needed.
    Returns the bottom y (Emu) after the last row drawn."""
    bg = bg or CHIP_BG
    fg = fg or _accent()
    x = left
    y = top
    row_h = Inches(0.3)
    pad_x = Inches(0.14)
    char_w = Emu(int(Inches(1) * (size / 11) * 0.075))
    for item in items:
        w = Emu(int(char_w) * len(item)) + pad_x * 2
        if x + w > left + max_width and x != left:
            x = left
            y += row_h + Inches(0.08)
        chip = slide.shapes.add_shape(5, x, y, w, row_h)
        try:
            chip.adjustments[0] = 0.5
        except (IndexError, AttributeError):
            pass
        chip.fill.solid()
        chip.fill.fore_color.rgb = bg
        chip.line.fill.background()
        chip.shadow.inherit = False
        tf = chip.text_frame
        tf.word_wrap = False
        tf.margin_left = pad_x
        tf.margin_right = pad_x
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = fg
        x += w + Inches(0.1)
    return y + row_h


def _icon_dot(slide, cx, cy, diameter, color):
    dot = slide.shapes.add_shape(9, cx, cy, diameter, diameter)
    _fill(dot, color)
    dot.shadow.inherit = False
    return dot


def _format_date_range(date_range: dict) -> str:
    """GA4 and Search Console are pulled with different windows (GSC lags
    ~3 days, see _gather_report_data) — when both are present and differ,
    state each explicitly rather than picking one and implying it covers
    both, since a reader comparing this slide's numbers against their own
    Search Console pull would otherwise get a mismatch with no explanation."""
    from datetime import date as _date

    def _fmt(iso: str) -> str:
        return _date.fromisoformat(iso).strftime("%b %d, %Y")

    ga4_start, ga4_end = date_range.get("ga4_start"), date_range.get("ga4_end")
    gsc_start, gsc_end = date_range.get("gsc_start"), date_range.get("gsc_end")
    parts = []
    if ga4_start and ga4_end:
        parts.append(f"Analytics: {_fmt(ga4_start)} – {_fmt(ga4_end)}")
    if gsc_start and gsc_end and (gsc_start, gsc_end) != (ga4_start, ga4_end):
        parts.append(f"Search Console: {_fmt(gsc_start)} – {_fmt(gsc_end)}")
    return "  ·  ".join(parts)


def _ga4_date_span(date_range: dict | None) -> str:
    """'Aug 01, 2026 – Aug 27, 2026' for the GA4 window in analytics["date_range"],
    or '' if unavailable — used so each Analytics slide states which dates its
    own numbers cover, not just the title slide."""
    if not date_range or not (date_range.get("ga4_start") and date_range.get("ga4_end")):
        return ""
    from datetime import date as _date

    fmt = lambda iso: _date.fromisoformat(iso).strftime("%b %d, %Y")
    return f"{fmt(date_range['ga4_start'])} – {fmt(date_range['ga4_end'])}"


def _gsc_date_span(date_range: dict | None) -> str:
    """Same as _ga4_date_span but for the Search Console window."""
    if not date_range or not (date_range.get("gsc_start") and date_range.get("gsc_end")):
        return ""
    from datetime import date as _date

    fmt = lambda iso: _date.fromisoformat(iso).strftime("%b %d, %Y")
    return f"{fmt(date_range['gsc_start'])} – {fmt(date_range['gsc_end'])}"


def add_title_slide(
    prs: Presentation, client_name: str, website_url: str = "", subtitle: str = "Web and SEO Audit",
    logo_bytes: bytes | None = None, analytics: dict | None = None,
):
    slide = _blank_slide(prs)
    _body_bg(slide)  # not used visually but keeps consistency; overwritten below
    top_bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.15))
    _fill(top_bar, DARK)
    top_bar.shadow.inherit = False

    white_area = slide.shapes.add_shape(1, 0, Inches(0.15), SLIDE_W, Inches(3.6))
    _fill(white_area, WHITE)
    white_area.shadow.inherit = False

    ring_box = (Inches(11.0), Inches(0.55), Inches(2.0), Inches(2.0))
    if logo_bytes:
        try:
            # Fit within the ring's box without distorting aspect ratio —
            # measure first, then scale by whichever dimension is larger.
            from PIL import Image

            img = Image.open(BytesIO(logo_bytes))
            w, h = img.size
            box_w, box_h = ring_box[2], ring_box[3]
            if w >= h:
                pic_w, pic_h = box_w, int(box_w * h / w)
            else:
                pic_h, pic_w = box_h, int(box_h * w / h)
            left = ring_box[0] + (box_w - pic_w) // 2
            top = ring_box[1] + (box_h - pic_h) // 2
            slide.shapes.add_picture(BytesIO(logo_bytes), left, top, width=pic_w, height=pic_h)
        except Exception:
            logo_bytes = None  # corrupt/unreadable image — fall through to the decorative ring
    if not logo_bytes:
        # decorative accent motif, echoes the score-ring circles used later in the deck
        ring = slide.shapes.add_shape(9, *ring_box)
        ring.fill.background()
        ring.line.color.rgb = _accent()
        ring.line.width = Pt(3)
        ring.shadow.inherit = False

    _textbox(
        slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
        client_name, size=44, bold=True, color=_accent(), align=PP_ALIGN.LEFT,
    )
    domain = website_url.replace("https://", "").replace("http://", "").rstrip("/")
    if domain:
        _textbox(slide, Inches(1.5), Inches(2.35), Inches(10), Inches(0.5), domain, size=16, color=TEXT_MUTED)

    blue_area = slide.shapes.add_shape(1, 0, Inches(3.75), SLIDE_W, Inches(3.75))
    _fill(blue_area, _accent())
    blue_area.shadow.inherit = False

    _textbox(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6), subtitle, size=24, bold=True, color=WHITE)

    from datetime import date

    date_range_text = _format_date_range(analytics["date_range"]) if analytics and analytics.get("date_range") else ""
    footer_text = f"{date.today().strftime('%B %Y')}   |   Data window — {date_range_text}" if date_range_text else date.today().strftime("%B %Y")
    _textbox(
        slide, Inches(1.5), Inches(5.15), Inches(10), Inches(0.4),
        footer_text, size=13, color=RGBColor(0xE0, 0xE0, 0xE0),
    )
    return slide


def add_section_slide(prs: Presentation, client_name: str, section_title: str):
    slide = _blank_slide(prs)
    top_bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.15))
    _fill(top_bar, DARK)
    top_bar.shadow.inherit = False

    white_area = slide.shapes.add_shape(1, 0, Inches(0.15), SLIDE_W, Inches(3.6))
    _fill(white_area, WHITE)
    white_area.shadow.inherit = False

    _textbox(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(1.0), client_name, size=36, bold=True, color=_accent())

    blue_area = slide.shapes.add_shape(1, 0, Inches(3.75), SLIDE_W, Inches(3.75))
    _fill(blue_area, _accent())
    blue_area.shadow.inherit = False
    _textbox(slide, Inches(1.5), Inches(4.6), Inches(10), Inches(0.6), section_title, size=24, bold=True, color=WHITE)
    return slide


def _score_color(score: int | None):
    if score is None:
        return TEXT_MUTED
    if score >= 90:
        return GOOD
    if score >= 50:
        return WARN
    return BAD


def _score_ring(slide, cx, cy, diameter, score: int | None, label: str):
    color = _score_color(score)
    # faint full backing ring so the colored arc reads as a gauge, not a plain circle
    backing = slide.shapes.add_shape(9, cx, cy, diameter, diameter)
    backing.fill.background()
    backing.line.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
    backing.line.width = Pt(7)
    backing.shadow.inherit = False

    ring = slide.shapes.add_shape(9, cx, cy, diameter, diameter)
    ring.fill.solid()
    ring.fill.fore_color.rgb = WHITE
    ring.line.color.rgb = color
    ring.line.width = Pt(6)
    ring.shadow.inherit = False
    tf = ring.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(score) if score is not None else "—"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = color
    _textbox(
        slide, cx - Inches(0.3), cy + diameter + Inches(0.05), diameter + Inches(0.6), Inches(0.35),
        label, size=12, color=TEXT_DARK, align=PP_ALIGN.CENTER,
    )


def _score_legend(slide, left, top):
    """Small colored-square + range legend, like the 0-49 / 50-89 / 90-100 chips in the sample."""
    items = [(BAD, "0–49"), (WARN, "50–89"), (GOOD, "90–100")]
    x = left
    for color, label in items:
        sq = slide.shapes.add_shape(1, x, top, Inches(0.14), Inches(0.14))
        _fill(sq, color)
        sq.shadow.inherit = False
        _textbox(slide, x + Inches(0.2), top - Inches(0.04), Inches(0.9), Inches(0.25), label, size=10, color=TEXT_MUTED)
        x += Inches(1.05)


def add_pagespeed_slide(prs: Presentation, mobile: dict | None, desktop: dict | None):
    slide = _blank_slide(prs)
    _content_header(slide, "Website Performance")
    _textbox(slide, Inches(9.5), Inches(0.3), Inches(3.3), Inches(0.4), "Source: Google PageSpeed Insights", size=11, color=TEXT_MUTED)

    for i, (label, result) in enumerate([("Mobile", mobile), ("Desktop", desktop)]):
        top = Inches(1.1) + Emu(i * Inches(2.9))
        card = _card(slide, Inches(0.4), top, Inches(9.0), Inches(2.5))
        _textbox(slide, Inches(0.6), top + Inches(0.1), Inches(3), Inches(0.4), label, size=16, bold=True)
        if not result:
            _textbox(slide, Inches(0.6), top + Inches(0.9), Inches(6), Inches(0.5), "Not run", size=13, color=TEXT_MUTED)
            continue
        scores = result.get("scores", {})
        labels = [("performance", "Performance"), ("seo", "SEO"), ("accessibility", "Accessibility"), ("best_practices", "Best Practices")]
        for j, (key, disp) in enumerate(labels):
            cx = Inches(0.9) + Emu(j * Inches(2.1))
            _score_ring(slide, cx, top + Inches(0.6), Inches(1.1), scores.get(key), disp)

    _score_legend(slide, Inches(9.8), Inches(1.6))
    _textbox(
        slide, Inches(9.8), Inches(2.0), Inches(3.1), Inches(1.3),
        "Best Practice: A score of 90+ is excellent. 50–89 is good but needs "
        "improvement. Below 50 is poor.",
        size=12, color=TEXT_DARK,
    )

    insights = []
    if mobile and desktop:
        m_scores, d_scores = mobile.get("scores", {}), desktop.get("scores", {})
        for key, disp in [("performance", "Performance"), ("seo", "SEO"), ("accessibility", "Accessibility"), ("best_practices", "Best Practices")]:
            m, d = m_scores.get(key), d_scores.get(key)
            if m is not None and m < 50:
                insights.append(f"Mobile {disp} is {m} — below the poor threshold, hurting mobile rankings directly.")
        if m_scores.get("performance") is not None and d_scores.get("performance") is not None:
            gap = d_scores["performance"] - m_scores["performance"]
            if gap > 15:
                insights.append(f"Mobile Performance trails desktop by {gap} points — most traffic is mobile, so this is the higher-impact fix.")
        worst_key, worst_val = None, 101
        for key, disp in [("performance", "Performance"), ("seo", "SEO"), ("accessibility", "Accessibility"), ("best_practices", "Best Practices")]:
            for scores in (m_scores, d_scores):
                v = scores.get(key)
                if v is not None and v < worst_val:
                    worst_val, worst_key = v, disp
        if worst_key and worst_val < 90:
            insights.append(f"Weakest area overall: {worst_key} at {worst_val}.")
    if insights:
        _insights_strip(slide, Inches(9.8), Inches(4.6), Inches(3.1), insights)
    return slide


def add_pagespeed_issues_slide(prs: Presentation, mobile: dict | None, desktop: dict | None):
    """The score-ring slide above says *how bad* PageSpeed is; this says
    *why* — the actual Lighthouse audit failures (render-blocking resources,
    unoptimized images, unused JS, etc), same ones PSI's own dashboard lists
    under "Opportunities"/"Diagnostics", instead of just a bare score."""
    rows = []
    for label, result in (("Mobile", mobile), ("Desktop", desktop)):
        for issue in (result or {}).get("issues", []):
            rows.append((issue["title"], label, issue.get("impact") or "—"))
    if not rows:
        return None

    rows.sort(key=lambda r: r[1])  # group by strategy so Mobile/Desktop aren't interleaved
    insights = []
    worst = max(
        (issue for result in (mobile, desktop) if result for issue in result.get("issues", [])),
        key=lambda i: i["savings_ms"], default=None,
    )
    if worst and worst["savings_ms"]:
        insights.append(f"Biggest opportunity: \"{worst['title']}\" — {worst['impact']}.")
    insights.append(f"{len(rows)} PageSpeed issue(s) found across Mobile/Desktop Lighthouse audits.")

    return _table_slide(
        prs, "Website Performance — Issue Details", ["Issue", "Where", "Impact"], rows,
        col_widths=[5.3, 1.3, 5.5], source="Google PageSpeed Insights", insights=insights,
    )


def _issue_row(slide, left, top, width, text, severity="warn"):
    color = BAD if severity == "error" else WARN
    dot = slide.shapes.add_shape(9, left, top + Inches(0.06), Inches(0.12), Inches(0.12))
    _fill(dot, color)
    dot.shadow.inherit = False
    _textbox(slide, left + Inches(0.25), top, width - Inches(0.25), Inches(0.35), text, size=13)


_CRAWLED_PAGE_CATEGORIES = ["Blocked", "Redirect", "Have issues", "Broken", "Healthy"]


def add_site_health_slide(prs: Presentation, audit: dict, site_audit_overview: dict | None = None):
    slide = _blank_slide(prs)
    _content_header(slide, "Understanding Current Scenario")
    if site_audit_overview and site_audit_overview.get("export_date"):
        _textbox(
            slide, Inches(8.3), Inches(0.3), Inches(4.5), Inches(0.4),
            f"Source: Semrush Site Audit ({site_audit_overview['export_date']})", size=11, color=TEXT_MUTED,
        )

    card1 = _card(slide, Inches(0.6), Inches(1.2), Inches(3.6), Inches(2.6))
    _textbox(slide, Inches(0.8), Inches(1.35), Inches(3), Inches(0.4), "Crawled Pages", size=15, bold=True)

    if site_audit_overview:
        health_pct = site_audit_overview.get("site_health_pct")
        _textbox(
            slide, Inches(0.8), Inches(1.68), Inches(2.5), Inches(0.5),
            str(site_audit_overview.get("pages_total", "—")), size=26, bold=True, color=_accent(),
        )
        category_colors = {
            "blocked": TEXT_MUTED,
            "redirect": RGBColor(0x5B, 0x5F, 0xE0),
            "have_issues": WARN,
            "broken": BAD,
            "healthy": GOOD,
        }
        y = Inches(2.18)
        for category in _CRAWLED_PAGE_CATEGORIES:
            key = category.lower().replace(" ", "_")
            count = site_audit_overview.get(f"{key}_count")
            pct = site_audit_overview.get(f"{key}_pct")
            if count is None:
                continue
            _icon_dot(slide, Inches(0.85), y + Inches(0.06), Inches(0.11), category_colors[key])
            _textbox(slide, Inches(1.05), y, Inches(2.9), Inches(0.28), f"{category}: {count} ({pct}%)", size=11.5, color=TEXT_DARK)
            y += Inches(0.29)
    else:
        # Crawled Pages / Site Health are Semrush-only now — no own-crawl
        # fallback. A client with no Site Audit Overview PDF uploaded yet
        # gets an explicit "no data" note instead of a silently-substituted
        # (and less accurate) own-crawl approximation.
        health_pct = None
        _textbox(
            slide, Inches(0.8), Inches(1.75), Inches(3.0), Inches(0.9),
            "No Semrush Site Audit data uploaded yet.", size=12.5, color=TEXT_MUTED,
        )

    card2 = _card(slide, Inches(4.5), Inches(1.2), Inches(3.6), Inches(2.6))
    _textbox(slide, Inches(4.7), Inches(1.35), Inches(3), Inches(0.4), "Site Health", size=15, bold=True)
    _score_ring(slide, Inches(5.6), Inches(1.85), Inches(1.4), health_pct, "full site crawl")
    if site_audit_overview and site_audit_overview.get("ai_search_health_pct") is not None:
        _textbox(
            slide, Inches(4.7), Inches(3.45), Inches(3.2), Inches(0.3),
            f"AI Search Health: {site_audit_overview['ai_search_health_pct']}%", size=11.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
        )

    card3 = _card(slide, Inches(8.4), Inches(1.2), Inches(4.3), Inches(2.6))
    _textbox(slide, Inches(8.6), Inches(1.35), Inches(4), Inches(0.4), "Foundations", size=15, bold=True)
    rows = [
        ("HTTPS", audit.get("https")),
        ("robots.txt present", audit.get("robots_txt", {}).get("present")),
        ("sitemap.xml present", audit.get("sitemap", {}).get("present")),
        ("Homepage reachable", audit.get("reachable")),
    ]
    for i, (label, ok) in enumerate(rows):
        y = Inches(1.85) + Emu(i * Inches(0.4))
        _icon_dot(slide, Inches(8.6), y + Inches(0.08), Inches(0.14), GOOD if ok else BAD)
        _textbox(slide, Inches(8.85), y, Inches(2.1), Inches(0.35), label, size=13)
        mark = "OK" if ok else "Missing"
        _textbox(slide, Inches(11.2), y, Inches(1.3), Inches(0.35), mark, size=13, bold=True, color=(GOOD if ok else BAD))

    _textbox(slide, Inches(0.6), Inches(4.1), Inches(9), Inches(0.35), f"Domain: {audit.get('url', '')}", size=13, color=TEXT_MUTED)
    return slide


def add_site_structure_slide(prs: Presentation, site_audit_pages_rows: list[dict] | None):
    """Directory-level rollup of the full-site crawl, styled after Semrush's
    own "Site Structure" widget: a root domain row (with its total URL
    count) on top, each top-level directory folder underneath with its own
    URL count, and — for directories that actually branch into further
    sub-directories (e.g. /ca, /product) — the top sub-directories nested
    beneath it, indented one level deeper, matching Semrush's expandable
    folder rows. Derived entirely from Semrush Site Audit's per-page export
    (page_url) we already parse for the SEO Issues / Tech Fixes slides — no
    new Semrush upload needed, just a grouping."""
    if not site_audit_pages_rows:
        return None

    # WordPress auto-generates a permalink for every post under one of these
    # base prefixes (post/page slugs, author archives, tag/category archives).
    # Each is an individual item or archive page, not a real content section
    # a client would recognize as a site "directory" — bucketing them
    # standalone produces phantom top-level directories (e.g. "/post" with
    # dozens of URLs) that don't correspond to anything in the client's nav.
    # Fold them into "/blog" alongside the section they actually belong to.
    _WP_ARCHIVE_PREFIXES = {"post", "author", "tag", "category", "page"}

    # A crawl can include a handful of pages from a subdomain (e.g. a
    # helpdesk/portal subdomain) alongside the main site — picking "any"
    # domain from a set is non-deterministic and previously surfaced a
    # rarely-crawled subdomain as the header instead of the actual site.
    # The domain with the most crawled pages is the real site, and every
    # other subdomain's paths are excluded below so they can't be
    # misattributed into this domain's directory structure.
    domain_counts = Counter(urlparse(r.get("page_url") or "").netloc for r in site_audit_pages_rows)
    domain_counts.pop("", None)
    if not domain_counts:
        return None
    domain = domain_counts.most_common(1)[0][0]

    top_counts: dict[str, int] = {}
    child_counts: dict[str, Counter] = {}
    root_count = 0
    for r in site_audit_pages_rows:
        parsed = urlparse(r.get("page_url") or "")
        if parsed.netloc != domain:
            continue
        segments = [s for s in parsed.path.split("/") if s]
        if not segments:
            # A bare "/" page — already covered by the domain row itself, so
            # it doesn't need its own "/ (root)" child row.
            root_count += 1
            continue
        if segments[0].lower() in _WP_ARCHIVE_PREFIXES and len(segments) > 1:
            directory = "/blog"
            sub_directory = f"/blog/{segments[0]}"
        else:
            directory = f"/{segments[0]}"
            sub_directory = f"{directory}/{segments[1]}" if len(segments) > 1 else None
        top_counts[directory] = top_counts.get(directory, 0) + 1
        if sub_directory:
            child_counts.setdefault(directory, Counter())[sub_directory] += 1

    ranked = sorted(top_counts.items(), key=lambda kv: -kv[1])
    domain_total = sum(top_counts.values()) + root_count

    slide = _blank_slide(prs)
    _content_header(slide, "Website Structure")
    _textbox(slide, Inches(8.3), Inches(0.3), Inches(4.5), Inches(0.4), "Source: Semrush Site Audit", size=11, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    # A real node-and-connector tree instead of an indented table — same
    # underlying domain -> directory -> sub-directory data as before, just
    # drawn as a hierarchy diagram (client asked for "hierarchy kind of
    # website structure"). Slide width caps this at 6 top-level directory
    # boxes and 2 sub-directory boxes per parent before it gets cramped;
    # anything beyond that is called out in the insight line below instead
    # of silently dropped.
    MAX_TOP = 6
    MAX_SUB = 2
    shown = ranked[:MAX_TOP]

    root_w, root_h = 3.0, 0.62
    root_left = 0.6 + (12.1 - root_w) / 2
    root_top = 1.05
    _tree_box(slide, Inches(root_left), Inches(root_top), Inches(root_w), Inches(root_h), domain, domain_total, bold=True, fill=HEADER_ROW_BG, name_size=13, count_size=11)

    n2 = max(len(shown), 1)
    gap2 = 0.15
    box2_w = (12.1 - (n2 - 1) * gap2) / n2
    box2_h = 0.68
    spine1_y = root_top + root_h + 0.22
    box2_top = spine1_y + 0.15

    root_center_x = root_left + root_w / 2
    _tree_line_v(slide, Inches(root_center_x), Inches(root_top + root_h), Inches(spine1_y))

    centers2 = [0.6 + box2_w / 2 + i * (box2_w + gap2) for i in range(n2)]
    if len(centers2) > 1:
        _tree_line_h(slide, Inches(centers2[0]), Inches(centers2[-1]), Inches(spine1_y))

    for i, (directory, count) in enumerate(shown):
        left2 = 0.6 + i * (box2_w + gap2)
        cx = centers2[i]
        _tree_line_v(slide, Inches(cx), Inches(spine1_y), Inches(box2_top))
        label = _truncate_cell(directory, box2_w - 0.2, size_pt=11)
        _tree_box(slide, Inches(left2), Inches(box2_top), Inches(box2_w), Inches(box2_h), label, count, name_size=11, count_size=9.5)

        subs = child_counts.get(directory)
        if not subs or len(subs) < 2:
            continue
        top_subs = sorted(subs.items(), key=lambda kv: -kv[1])[:MAX_SUB]
        n3 = len(top_subs)
        gap3 = 0.1
        box3_w = (box2_w - (n3 - 1) * gap3) / n3
        box3_h = 0.58
        spine2_y = box2_top + box2_h + 0.16
        box3_top = spine2_y + 0.12
        parent_bottom = box2_top + box2_h
        _tree_line_v(slide, Inches(cx), Inches(parent_bottom), Inches(spine2_y))
        centers3 = [left2 + box3_w / 2 + j * (box3_w + gap3) for j in range(n3)]
        if len(centers3) > 1:
            _tree_line_h(slide, Inches(centers3[0]), Inches(centers3[-1]), Inches(spine2_y))
        for j, (sub_directory, sub_count) in enumerate(top_subs):
            left3 = left2 + j * (box3_w + gap3)
            _tree_line_v(slide, Inches(centers3[j]), Inches(spine2_y), Inches(box3_top))
            sub_label = _truncate_cell(sub_directory.split("/")[-1], box3_w - 0.15, size_pt=9.5)
            _tree_box(slide, Inches(left3), Inches(box3_top), Inches(box3_w), Inches(box3_h), f"/{sub_label}", sub_count, name_size=9.5, count_size=8.5)

    insights = [f"{domain_total:,} total crawled URLs across {len(ranked)} top-level director{'y' if len(ranked) == 1 else 'ies'}."]
    if len(ranked) > MAX_TOP:
        insights.append(f"Showing the top {MAX_TOP} directories by URL count — {len(ranked) - MAX_TOP} more not shown here.")
    _insights_strip(slide, Inches(0.6), Inches(5.4), Inches(11.9), insights)
    return slide


def _tree_box(slide, left, top, width, height, name: str, count: int, bold: bool = False, fill=None, name_size: float = 11, count_size: float = 9.5):
    """One node in the Website Structure hierarchy diagram — a compact
    rounded box with the directory name on top and its URL count below,
    matching the deck's existing card styling (see _card) rather than a
    plain shape, so the tree reads as part of the same visual system."""
    box = slide.shapes.add_shape(5, left, top, width, height)  # rounded rectangle
    try:
        box.adjustments[0] = 0.12
    except (IndexError, AttributeError):
        pass
    box.fill.solid()
    box.fill.fore_color.rgb = fill or WHITE
    box.line.color.rgb = _accent() if bold else CARD_BORDER
    box.line.width = Pt(1.25 if bold else 0.75)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(2)
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    r0 = p0.add_run()
    r0.text = name
    r0.font.size = Pt(name_size)
    r0.font.bold = bold
    r0.font.color.rgb = TEXT_DARK
    p1 = tf.add_paragraph()
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = f"{count:,} URLs"
    r1.font.size = Pt(count_size)
    r1.font.bold = True
    r1.font.color.rgb = _accent() if bold else TEXT_MUTED
    return box


def _tree_line_v(slide, x, y1, y2):
    line = slide.shapes.add_shape(1, x, min(y1, y2), Pt(1.5), abs(int(y2) - int(y1)) or Pt(1))
    _fill(line, CARD_BORDER)
    line.shadow.inherit = False


def _tree_line_h(slide, x1, x2, y):
    line = slide.shapes.add_shape(1, min(x1, x2), y, abs(int(x2) - int(x1)) or Pt(1), Pt(1.5))
    _fill(line, CARD_BORDER)
    line.shadow.inherit = False


def add_company_overview_slide(prs: Presentation, client_name: str, summary: str):
    slide = _blank_slide(prs)
    _content_header(slide, "Company Overview")
    _card(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(3.4))
    _textbox(slide, Inches(0.9), Inches(1.45), Inches(11.5), Inches(0.4), client_name, size=17, bold=True)
    _textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.4), summary, size=14)
    return slide


def _icp_block(slide, x, top, width, label, items, size=11):
    """One labeled chip-row block in the Ideal Customer Profile column.
    Returns the bottom y (Emu) after the block."""
    if not items:
        return top
    _textbox(slide, x, top, width, Inches(0.24), label.upper(), size=9.5, bold=True, color=TEXT_MUTED)
    bottom = _chip_row(slide, x, top + Inches(0.26), items, width, size=size)
    return bottom + Inches(0.14)


def add_company_overview_extracted_slide(prs: Presentation, client_name: str, overview: dict):
    """Renders the Gemini-extracted About/Products/KPIs/ICP overview as two
    dense columns — narrative + KPIs on the left, Ideal Customer Profile as
    wrapping chip rows on the right, contact/registration pinned to the
    footer. Chips pack more into less vertical space than a bullet list."""
    slide = _blank_slide(prs)
    _content_header(slide, overview.get("company_name") or client_name, eyebrow="Company Overview")
    _textbox(slide, Inches(9.6), Inches(0.4), Inches(3.3), Inches(0.3), "Source: site crawl + Gemini", size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)

    top = Inches(1.15)
    height = Inches(5.55)
    left_width = Inches(5.55)
    gap = Inches(0.35)
    right_x = Inches(0.6) + left_width + gap
    right_width = Inches(12.1) - left_width - gap

    _card(slide, Inches(0.6), top, left_width, height)
    _card(slide, right_x, top, right_width, height)

    pad = Inches(0.25)
    y = top + pad

    description = overview.get("description")
    if description:
        _textbox(slide, Inches(0.6) + pad, y, left_width - pad * 2, Inches(1.6), description, size=12.5)
        chars_per_line = 50
        lines = max(1, -(-len(description) // chars_per_line))
        y += Emu(lines * Inches(0.22)) + Inches(0.22)

    kpis = overview.get("kpis") or []
    if kpis:
        _textbox(slide, Inches(0.6) + pad, y, left_width - pad * 2, Inches(0.26), "KEY PERFORMANCE INDICATORS", size=9.5, bold=True, color=_accent())
        y += Inches(0.3)
        kpi_width = left_width - pad * 2 - Inches(0.18)
        line_h = Inches(11.5 * 0.02)
        for k in kpis[:6]:
            lines = _wrap_lines(k, kpi_width, size_pt=11.5)
            _icon_dot(slide, Inches(0.6) + pad, y + Inches(0.07), Inches(0.08), _accent())
            _textbox(slide, Inches(0.6) + pad + Inches(0.18), y, kpi_width, line_h * lines, k, size=11.5)
            y += line_h * lines + Inches(0.06)
        y += Inches(0.1)

    industries = overview.get("industries") or []
    if industries:
        _textbox(slide, Inches(0.6) + pad, y, left_width - pad * 2, Inches(0.24), "INDUSTRIES SERVED", size=9.5, bold=True, color=_accent())
        y = _chip_row(slide, Inches(0.6) + pad, y + Inches(0.26), industries[:8], left_width - pad * 2, size=10.5, bg=ROW_ALT, fg=TEXT_DARK) + Inches(0.1)

    ry = top + pad
    rx = right_x + pad
    rw = right_width - pad * 2
    _textbox(slide, rx, ry, rw, Inches(0.28), "IDEAL CUSTOMER PROFILE", size=10.5, bold=True, color=_accent())
    ry += Inches(0.36)

    target_country = overview.get("target_country")
    target_market = overview.get("target_market")
    if target_country or target_market:
        half = rw / 2 - Inches(0.1)
        if target_country:
            _textbox(slide, rx, ry, half, Inches(0.24), "TARGET COUNTRY", size=9, bold=True, color=TEXT_MUTED)
            _textbox(slide, rx, ry + Inches(0.22), half, Inches(0.3), target_country, size=12.5, bold=True)
        if target_market:
            _textbox(slide, rx + half + Inches(0.2), ry, half, Inches(0.24), "TARGET", size=9, bold=True, color=TEXT_MUTED)
            _textbox(slide, rx + half + Inches(0.2), ry + Inches(0.22), half, Inches(0.3), target_market, size=12.5, bold=True)
        ry += Inches(0.62)

    ry = _icp_block(slide, rx, ry, rw, "Primary buyers", (overview.get("primary_buyers") or [])[:6])
    ry = _icp_block(slide, rx, ry, rw, "Daily users", (overview.get("daily_users") or [])[:6])
    ry = _icp_block(slide, rx, ry, rw, "Beneficiaries", (overview.get("beneficiaries") or [])[:6])

    reg = overview.get("registration_info")
    contact = overview.get("contact")
    footer_bits = [b for b in [reg, contact] if b]
    if footer_bits:
        _textbox(slide, Inches(0.6), top + height + Inches(0.12), Inches(12.1), Inches(0.3), "  |  ".join(footer_bits), size=10.5, color=TEXT_MUTED)
    return slide


def add_solutions_products_slide(prs: Presentation, overview: dict):
    """Renders Solutions / Products-by-category / Industries — matches the
    agency sample deck's 'Solutions, Products & Industries' slide."""
    solutions = overview.get("solutions") or []
    products_by_category = overview.get("products_by_category") or {}
    industries = overview.get("industries") or []
    products_flat = overview.get("products") or []
    if not products_by_category and products_flat:
        products_by_category = {"Products & Services": products_flat}
    if not solutions and not products_by_category and not industries:
        return None

    slide = _blank_slide(prs)
    _content_header(slide, "Products & Services" if not solutions and not industries else "Solutions, Products & Industries")
    _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.9))
    y = Inches(1.35)

    if solutions:
        _textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.35), "Solutions", size=15, bold=True, color=_accent())
        y += Inches(0.4)
        for s in solutions[:8]:
            if y > Inches(6.6):
                break
            text = f"•  {s}"
            lines = _wrap_lines(text, Inches(11.1), size_pt=12)
            line_h = Inches(12 * 0.02)
            box = slide.shapes.add_textbox(Inches(1.1), y, Inches(11.1), line_h * lines)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(12)
            run.font.color.rgb = TEXT_DARK
            y += line_h * lines + Inches(0.06)
        y += Inches(0.15)

    if products_by_category:
        _textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.35), "Products (by category)", size=15, bold=True, color=_accent())
        y += Inches(0.4)
        for category, items in list(products_by_category.items())[:6]:
            if y > Inches(6.6):
                break
            _textbox(slide, Inches(0.9), y, Inches(11.1), Inches(0.3), category, size=13, bold=True)
            y += Inches(0.32)
            joined = ", ".join(items)
            lines = _wrap_lines(joined, Inches(11.3), size_pt=11)
            line_h = Inches(11 * 0.02)
            box = slide.shapes.add_textbox(Inches(1.0), y, Inches(11.3), line_h * lines)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = joined
            run.font.size = Pt(11)
            run.font.color.rgb = TEXT_MUTED
            y += line_h * lines + Inches(0.2)

    if industries:
        if y < Inches(6.7):
            _textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.35), "Industries", size=15, bold=True, color=_accent())
            y += Inches(0.4)
            _textbox(slide, Inches(0.9), y, Inches(11.3), Inches(0.5), ", ".join(industries[:12]), size=12)
        else:
            # Didn't fit on the main slide — a long Solutions/Products list
            # can push y past the card before Industries even starts (seen
            # live on a real report: Industries vanished with no trace).
            # Real client data never gets silently dropped for lack of
            # room — give it its own slide instead.
            overflow_slide = _blank_slide(prs)
            _content_header(overflow_slide, "Industries Served")
            _card(overflow_slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.6))
            _textbox(overflow_slide, Inches(0.9), Inches(1.4), Inches(11.3), Inches(1.0), ", ".join(industries[:12]), size=14)

    return slide


def add_tech_stack_slide(prs: Presentation, tech_stack: dict):
    """Detected CMS/framework/hosting/CDN/analytics — from response headers,
    DNS/PTR, and HTML markers. No credentials involved."""
    slide = _blank_slide(prs)
    _content_header(slide, "Tech Stack & Hosting")
    _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.6))
    y = Inches(1.35)

    facts = [
        ("Hostname", tech_stack.get("hostname")),
        ("IP address", tech_stack.get("ip")),
        ("Reverse DNS", tech_stack.get("reverse_dns")),
        ("HTTPS", "Yes" if tech_stack.get("https") else "No"),
    ]
    for label, value in facts:
        if not value:
            continue
        _textbox(slide, Inches(0.9), y, Inches(2.2), Inches(0.3), label, size=11, color=TEXT_MUTED)
        _textbox(slide, Inches(3.2), y, Inches(9.0), Inches(0.3), str(value), size=12, bold=True)
        y += Inches(0.34)

    y += Inches(0.2)
    detected = tech_stack.get("detected") or []
    if detected:
        _textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.35), "Detected technologies", size=14, bold=True, color=_accent())
        y += Inches(0.42)
        by_category: dict[str, list[str]] = {}
        for item in detected:
            by_category.setdefault(item["category"], []).append(item["name"])
        for category, names in by_category.items():
            if y > Inches(6.5):
                break
            _textbox(slide, Inches(0.9), y, Inches(2.4), Inches(0.3), category.title(), size=12, bold=True)
            _textbox(slide, Inches(3.2), y, Inches(9.0), Inches(0.3), ", ".join(names), size=12, color=TEXT_DARK)
            y += Inches(0.36)
    return slide


def add_seo_issues_slide(prs: Presentation, audit: dict, page_audit: dict | None, site_audit_issues: list[dict] | None = None):
    slide = _blank_slide(prs)
    _content_header(slide, "SEO Issues")

    if site_audit_issues:
        # Semrush Site Audit's own issue-type rollup — a real full-site crawl
        # result (hundreds of pages, ~95 issue categories), strictly richer
        # than our own homepage + 20-page checks below. Prefer it when
        # uploaded. Rows with 0 failed checks are noise (the issue TYPE was
        # checked for but never triggered) — drop them so real problems
        # aren't crowded out by "X (0 pages)" lines.
        nonzero = [r for r in site_audit_issues if (r.get("failed_checks") or 0) > 0]
        ranked = sorted(nonzero, key=lambda r: r.get("failed_checks") or 0, reverse=True)
        errors, warnings = [], []
        for row in ranked:
            label = f"{row.get('issue', 'Issue')} ({row.get('failed_checks', 0)} pages)"
            if str(row.get("issue_type", "")).strip().upper() == "ERROR":
                errors.append(label)
            else:
                warnings.append(label)
    else:
        issues = list(audit.get("issues", []))
        if page_audit:
            for page in page_audit.get("pages", []):
                for issue in page.get("issues", []):
                    issues.append(f"{issue} — {page['url']}")
        errors = [i for i in issues if any(k in i.lower() for k in ["not reachable", "https", "robots", "sitemap"])]
        warnings = [i for i in issues if i not in errors]

    if not errors and not warnings:
        card = _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.6))
        _textbox(slide, Inches(0.9), Inches(1.3), Inches(10), Inches(0.4), "No issues found on the checked pages.", size=14, color=GOOD)
        return slide

    # Errors and Warnings each get their own full-height container, side by
    # side, rather than stacking in one shared card — stacked lists with 10
    # items per group could run past the bottom of the slide.
    col_top, col_height = Inches(1.1), Inches(5.6)
    col_width = Inches(5.85)
    columns = [("Errors", errors, BAD, Inches(0.6)), ("Warnings", warnings, WARN, Inches(6.85))]
    row_h = Inches(0.32)
    max_rows = 10

    for group_label, group_issues, color, col_left in columns:
        if not group_issues:
            continue
        shown = group_issues[:max_rows]
        _card(slide, col_left, col_top, col_width, col_height)
        y = col_top + Inches(0.2)
        _textbox(slide, col_left + Inches(0.3), y, col_width - Inches(0.6), Inches(0.35), f"{group_label} ({len(shown)})", size=14, bold=True, color=color)
        rule = slide.shapes.add_shape(1, col_left + Inches(0.3), y + Inches(0.36), col_width - Inches(0.6), Pt(1.5))
        _fill(rule, color)
        rule.shadow.inherit = False
        y += Inches(0.55)
        for issue in shown:
            _issue_row(slide, col_left + Inches(0.3), y, col_width - Inches(0.6), issue, severity=("error" if color == BAD else "warn"))
            y += row_h
    return slide


def add_critical_issues_slide(
    prs: Presentation,
    site_audit_issues: list[dict] | None,
    site_audit_pages_rows: list[dict] | None = None,
    analytics: dict | None = None,
):
    """Standalone Critical Issues slide — ERROR-severity items only from
    Semrush Site Audit's issue rollup (issue_type == "ERROR"), matching the
    manual report's dedicated Critical Issues page. Distinct from the
    combined SEO Issues slide above, which mixes Errors + Warnings in one
    2-column layout — this is Errors alone, full width, so the most severe
    findings aren't sharing space with lower-priority warnings."""
    if not site_audit_issues:
        return None
    errors = [
        r for r in site_audit_issues
        if str(r.get("issue_type", "")).strip().upper() == "ERROR" and (r.get("failed_checks") or 0) > 0
    ]
    if not errors:
        return None

    ranked = sorted(errors, key=lambda r: r.get("failed_checks") or 0, reverse=True)
    rows = [
        (r.get("issue", "Issue"), f"{r.get('failed_checks', 0):,}", f"{r.get('total_checks', 0):,}" if r.get("total_checks") else "—")
        for r in ranked
    ]
    total_affected = sum(r.get("failed_checks") or 0 for r in errors)
    top = ranked[0]
    insights = [
        f"{len(errors)} distinct critical error type(s), {total_affected:,} total affected page-checks.",
        f"Most widespread: \"{top.get('issue')}\" — {top.get('failed_checks', 0):,} pages affected.",
    ]

    # Business-impact cross-reference: real GA4 pageviews that landed on
    # pages which are now broken/non-200, per the same Site Audit per-page
    # export the Website Structure slide uses. Only added when BOTH real
    # datasets are present and actually overlap — never estimated.
    if site_audit_pages_rows and analytics:
        broken_paths = {
            urlparse(r.get("page_url") or "").path.rstrip("/")
            for r in site_audit_pages_rows
            if str(r.get("http_status_code", "200")).strip() not in ("200", "")
        }
        top_pages = (analytics.get("top_pages") or {}).get("rows", [])
        matched_views = sum(
            int(float(p.get("page_views", 0) or 0))
            for p in top_pages
            if (p.get("path") or "").rstrip("/") in broken_paths
        )
        if matched_views > 0:
            insights.append(
                f"These broken/non-200 pages received {matched_views:,} real pageviews in the reporting window — "
                "that traffic (and any ranking credit) is being lost right now."
            )

    return _table_slide(
        prs, "Critical Issues", ["Issue", "Affected Pages", "Total Checked"], rows,
        col_widths=[7.0, 2.5, 2.6], source="Semrush Site Audit", insights=insights,
    )


def add_priority_issues_slide(
    prs: Presentation,
    site_audit_pages_rows: list[dict] | None,
    page_audit: dict | None,
    analytics: dict | None,
):
    """Ranks per-URL issues by real traffic impact instead of raw issue
    count: pageviews (GA4) + clicks (GSC), clicks weighted 3x since a lost
    search click is lost demand, a lost pageview could be internal/referral
    traffic that finds the page another way regardless of its search
    ranking. Needs at least one real per-URL issue source AND GA4 or GSC
    data — silently absent otherwise, same as the rest of this file's
    cross-referenced insights."""
    if not analytics:
        return None

    url_issues: list[tuple[str, str]] = []
    if site_audit_pages_rows:
        # GA4/GSC data below is matched by path alone (pagePath has no host
        # dimension) — a page on a different subdomain (agents., services.,
        # timesheet., unionagent.lumberfi.com, etc.) that happens to share a
        # path with the main site's homepage would otherwise silently
        # inherit the homepage's real pageviews/clicks. Only the dominant
        # crawled domain actually has per-path analytics behind it, so
        # every other host's issue rows are dropped here rather than scored
        # against traffic that was never theirs.
        domain_counts = Counter(urlparse(r.get("page_url") or "").netloc for r in site_audit_pages_rows)
        domain_counts.pop("", None)
        own_domain = domain_counts.most_common(1)[0][0] if domain_counts else None
        for r in site_audit_pages_rows:
            issues = r.get("issues")
            page_url = r.get("page_url")
            if not (issues and page_url):
                continue
            if own_domain and urlparse(page_url).netloc not in ("", own_domain):
                continue
            url_issues.append((page_url, str(issues)))
    elif page_audit:
        for p in page_audit.get("pages") or []:
            if p.get("issues") and p.get("url"):
                url_issues.append((p["url"], ", ".join(p["issues"])))
    if not url_issues:
        return None

    # Dedupe by normalized path — the homepage (or any page) can appear
    # twice as a trailing-slash variant ("https://site.com" vs.
    # "https://site.com/") when a Semrush Site Audit export lists both,
    # and both resolve to the same GA4/GSC path below — without this they'd
    # score identically and print as two adjacent, seemingly-repeated rows
    # on the same slide. First occurrence wins (order from the source data
    # is otherwise arbitrary here, before the traffic-based sort below).
    deduped_url_issues: dict[str, tuple[str, str]] = {}
    for url, issues in url_issues:
        key = urlparse(url).path.rstrip("/") or "/"
        if key not in deduped_url_issues:
            deduped_url_issues[key] = (url, issues)
    url_issues = list(deduped_url_issues.values())

    pageviews_by_path = {
        (p.get("path") or "").rstrip("/"): int(float(p.get("page_views", 0) or 0))
        for p in (analytics.get("top_pages") or {}).get("rows", [])
    }
    clicks_by_path = {
        urlparse(r.get("page") or "").path.rstrip("/"): int(r.get("clicks", 0) or 0)
        for r in (analytics.get("page_clicks") or {}).get("rows", [])
    }
    if not pageviews_by_path and not clicks_by_path:
        return None

    ranked = []
    for url, issues in url_issues:
        path = urlparse(url).path.rstrip("/")
        pageviews = pageviews_by_path.get(path, 0)
        clicks = clicks_by_path.get(path, 0)
        score = pageviews + clicks * 3
        if score > 0:
            ranked.append((url, issues, pageviews, clicks, score))
    if not ranked:
        return None

    ranked.sort(key=lambda r: r[4], reverse=True)
    headers = ["Page URL", "Issue(s)", "Pageviews", "GSC Clicks", "Priority Score"]
    col_widths = [3.6, 4.2, 1.3, 1.3, 1.7]
    shown = ranked[:14]
    rows = [
        (_truncate_cell(url, col_widths[0]), _truncate_cell(issues, col_widths[1]), f"{pv:,}", f"{cl:,}", f"{score:,}")
        for url, issues, pv, cl, score in shown
    ]
    top = ranked[0]
    insights = [
        f"{len(ranked)} issue-affected page(s) cross-referenced against real traffic — ranked by pageviews + clicks×3 (a lost search click is lost demand, weighted heavier than a pageview from another channel).",
        f"Highest priority: {top[0]} — {top[2]:,} pageviews, {top[3]:,} search clicks in the reporting window, still carrying: {top[1]}",
    ]
    slide = _table_slide(
        prs, "Priority Issues (by Traffic Impact)", headers, rows, col_widths=col_widths,
        source="Own crawl + Google Analytics + Search Console", insights=insights,
    )
    # Page URL cell links out to the live page itself — the full issue list
    # is already in the adjacent column (just truncated to fit), so this
    # link is for seeing the issue in context on the actual page, not for
    # more issue detail we don't otherwise have (no Semrush deep-link to
    # reuse here).
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        for i, (url, *_rest) in enumerate(shown, start=1):
            for run in table.cell(i, 0).text_frame.paragraphs[0].runs:
                run.hyperlink.address = url
        break
    return slide


# Semrush Site Audit's per-page Structured Data export has ~29 schema.org
# rich-result-type columns — this is the curated, business-relevant subset
# actually surfaced on the slide (nothing is lost on import, see
# STRUCTURED_DATA_COLUMN_ALIASES in semrush_parser.py; this is a display
# choice, not a parsing limit). Matches the style of real manual-report
# findings, e.g. "No Product schema (JSON-LD) on the 5 product pages"
# (confirmed from an EJTOY audit).
_STRUCTURED_DATA_TYPES = [
    ("Article", "article_items", "improves how blog/news content can appear in search"),
    ("FAQ", "faq_items", "enables expandable FAQ dropdowns directly in search results"),
    ("Product", "product_items", "enables price/availability rich results on product pages"),
    ("Review", "review_items", "enables star-rating rich results"),
    ("Local Business", "local_business_items", "enables map/business-info rich results"),
    ("How-to", "howto_items", "enables step-by-step rich results"),
    ("Breadcrumb", "breadcrumb_items", "shows the page's site hierarchy in search results"),
    ("Job Posting", "job_posting_items", "enables Google's dedicated job-search rich results"),
    ("Event", "event_items", "enables event date/venue rich results"),
]


# Every schema.org item-type field Semrush's Structured Data export tracks
# (STRUCTURED_DATA_COLUMN_ALIASES in semrush_parser.py) — a superset of the
# 9 curated rich-result types in _STRUCTURED_DATA_TYPES. Checked so a type
# no one asked to curate (Logo, Organization sitelinks, etc.) still shows up
# on the slide if it actually has real coverage, instead of silently having
# no row at all.
_ALL_SCHEMA_ITEM_FIELDS = [
    ("Article", "article_items"), ("Book", "book_items"), ("Breadcrumb", "breadcrumb_items"),
    ("Carousel", "carousel_items"), ("Course", "course_items"), ("Dataset", "dataset_items"),
    ("Employer Rating", "employer_rating_items"), ("Estimated Salary", "estimated_salary_items"),
    ("Event", "event_items"), ("Fact Check", "fact_check_items"), ("FAQ", "faq_items"),
    ("Guided Recipe", "guided_recipe_items"), ("How-to", "howto_items"), ("Job Posting", "job_posting_items"),
    ("Local Business", "local_business_items"), ("Logo", "logo_items"), ("Merchant Listing", "merchant_listing_items"),
    ("Movie", "movie_items"), ("Product", "product_items"), ("Product Group", "product_group_items"),
    ("Q&A", "qa_items"), ("Recipe", "recipe_items"), ("Review", "review_items"),
    ("Sitelinks Search Box", "sitelinks_searchbox_items"), ("Site Names", "site_names_items"),
    ("Software App", "software_app_items"), ("Vehicle Listing", "vehicle_listing_items"), ("Video", "video_items"),
]

# Contextual schema relevance: a page has to actually BE the shape a schema
# type is for before its absence counts as a gap — a SaaS site with zero
# careers pages doesn't need Job Posting schema, and flagging it as missing
# "on all 1,336 pages" is noise, not a finding. Product/blog/jobs/local/
# event are detected by URL pattern (cheap, no extra crawl); FAQ/how-to
# genuinely can't be — a FAQ section can live on any URL, including a
# homepage or product page — so those fall back to whatever text is
# available (URL slug + page title, when Site Audit Pages data is joined
# in) as a best-effort proxy. True content-scanning (question-shaped
# headings, numbered steps) would need our own crawler to fetch full page
# HTML — that infra existed once (see crawl_extras_2026-09-01) and was
# reverted for being too slow; this stays URL/title-only until that's
# safely rebuilt.
_PRODUCT_SHAPE_RE = re.compile(r"/(?:products?|shop|store|items?)/", re.IGNORECASE)
_BLOG_SHAPE_RE = re.compile(r"/(?:blog|news|articles?|posts?)/", re.IGNORECASE)
_JOBS_SHAPE_RE = re.compile(r"/(?:careers?|jobs?)/", re.IGNORECASE)
_LOCAL_SHAPE_RE = re.compile(r"/(?:locations?|store-locator|near-me|branch(?:es)?)/", re.IGNORECASE)
_EVENT_SHAPE_RE = re.compile(r"/events?/", re.IGNORECASE)
_FAQ_SHAPE_RE = re.compile(r"faq|frequently[\s-]asked[\s-]questions", re.IGNORECASE)
_HOWTO_SHAPE_RE = re.compile(r"how[\s-]to|step[\s-]by[\s-]step|\btutorial\b|\bguide\b", re.IGNORECASE)


def _classify_page_shapes(page_url: str, page_title: str = "") -> set[str]:
    path = urlparse(page_url or "").path
    shapes = set()
    if _PRODUCT_SHAPE_RE.search(path):
        shapes.add("product")
    if _BLOG_SHAPE_RE.search(path):
        shapes.add("blog")
    if _JOBS_SHAPE_RE.search(path):
        shapes.add("jobs")
    if _LOCAL_SHAPE_RE.search(path):
        shapes.add("local")
    if _EVENT_SHAPE_RE.search(path):
        shapes.add("event")
    text = f"{path} {page_title or ''}"
    if _FAQ_SHAPE_RE.search(text):
        shapes.add("faq")
    if _HOWTO_SHAPE_RE.search(text):
        shapes.add("howto")
    return shapes


# Which page shape(s) make a curated schema type's absence a real finding.
# A type with no entry here (Breadcrumb) applies site-wide — every page can
# carry one, so its denominator stays the full crawl, same as before.
_SCHEMA_TYPE_REQUIRED_SHAPES = {
    "article_items": {"blog"},
    "faq_items": {"faq"},
    "product_items": {"product"},
    "review_items": {"product"},  # reviews live on/near the product they review
    "local_business_items": {"local"},
    "howto_items": {"howto"},
    "job_posting_items": {"jobs"},
    "event_items": {"event"},
}
_SCHEMA_TYPE_SHAPE_NOUN = {
    "article_items": "blog/news",
    "faq_items": "FAQ-worthy",
    "product_items": "product",
    "review_items": "product",
    "local_business_items": "location",
    "howto_items": "how-to/guide",
    "job_posting_items": "job listing",
    "event_items": "event",
}


def add_structured_data_slide(
    prs: Presentation, structured_data_rows: list[dict], site_audit_pages_rows: list[dict] | None = None
):
    """Total URLs crawled vs. how many have any schema markup implemented,
    plus a per-type breakdown, worst-coverage first. Shows the 9 curated
    rich-result types PLUS any other tracked type (of all ~27 Semrush
    exports) that actually has real coverage — so if the pages missing a
    curated type (e.g. Product) carry some other schema instead, that shows
    up as its own row here rather than needing a separate slide to explain
    where the remaining % went. total_pages is the real crawl total
    (semrush_parser.py no longer caps structured_data rows at 500 — a real
    client site can have 1200+ crawled pages, and capping silently
    understated both this total and every coverage %).

    The coverage TABLE stays site-wide (X of ALL crawled pages) deliberately
    — narrowing its denominator to "relevant" pages would also silently
    exclude a genuinely valuable finding like a blog post carrying Product
    schema instead of Article (confirmed on a real Lumber crawl). Only the
    "No X schema found" INSIGHT bullets get contextual: a type mapped to a
    page shape (see _SCHEMA_TYPE_REQUIRED_SHAPES) is only flagged as missing
    when at least one page of that shape actually exists on the site — a
    SaaS site with zero careers pages doesn't need a "no Job Posting schema"
    bullet. site_audit_pages_rows (optional) supplies page_title for the
    FAQ/how-to shape check, which can't rely on URL alone."""
    total_pages = len(structured_data_rows)
    if not total_pages:
        return None

    title_by_url: dict[str, str] = {}
    if site_audit_pages_rows:
        for r in site_audit_pages_rows:
            url = r.get("page_url")
            if url:
                title_by_url[url] = r.get("page_title") or ""
    row_shapes = [
        _classify_page_shapes(r.get("page_url") or "", title_by_url.get(r.get("page_url") or "", ""))
        for r in structured_data_rows
    ]
    relevant_counts = {
        field: sum(1 for shapes in row_shapes if shapes & required)
        for field, required in _SCHEMA_TYPE_REQUIRED_SHAPES.items()
    }

    pages_with_any_schema = sum(
        1 for r in structured_data_rows
        if _num(r.get("schema_jsonld")) > 0 or _num(r.get("schema_microdata")) > 0
    )
    any_schema_pct = 100 * pages_with_any_schema / total_pages

    curated_fields = {field for _label, field, _benefit in _STRUCTURED_DATA_TYPES}
    coverage = [
        (label, field, benefit, sum(1 for r in structured_data_rows if _num(r.get(field)) > 0))
        for label, field, benefit in _STRUCTURED_DATA_TYPES
    ]
    # Any non-curated type that actually has coverage earns its own row too
    # — otherwise a client whose gap pages use e.g. Logo schema instead of
    # Product would show 85% Product and nothing explaining the other 15%.
    extra_coverage = [
        (label, field, None, sum(1 for r in structured_data_rows if _num(r.get(field)) > 0))
        for label, field in _ALL_SCHEMA_ITEM_FIELDS
        if field not in curated_fields
    ]
    coverage += [c for c in extra_coverage if c[3] > 0]
    # Highest-coverage first — with an extra non-curated row possibly added
    # above, the table's row cap (9 when insights are shown) must never
    # silently truncate the row(s) that actually have real coverage in
    # favor of interchangeable 0% rows; the "missing" story is already
    # covered by the insight bullets below regardless of which 0% rows
    # make the cut.
    coverage.sort(key=lambda c: c[3], reverse=True)

    headers = ["Schema Type", "Pages With It", "Coverage"]
    col_widths = [4.0, 2.5, 2.5]
    rows = [
        (label, f"{pages_with:,} / {total_pages:,}", f"{100 * pages_with / total_pages:.0f}%")
        for label, _field, _benefit, pages_with in coverage
    ]

    # The per-type coverage rows above are NOT mutually exclusive (a page can
    # carry Product AND Review AND nothing-else-tracked at once), so they
    # never sum to 100% and were never meant to — confirmed this reads as
    # "where's the rest of the percentage?" to a non-technical reviewer.
    # The fix: put the complementary "and here's the rest" stat (pages with
    # NO tracked schema at all) as the SECOND insight, right after the
    # headline stat — guaranteed to render even when the per-type "missing
    # X" bullets after it get cut for space — so the two numbers that
    # actually do sum to 100% of crawled pages are always shown together.
    insights = [
        f"{pages_with_any_schema:,} of {total_pages:,} crawled URLs ({any_schema_pct:.0f}%) have structured data "
        "(schema.org JSON-LD or Microdata) implemented.",
    ]
    best = max(coverage, key=lambda c: c[3])
    gap = total_pages - pages_with_any_schema
    if best[3] > 0 and gap > 0:
        insights.append(
            f"The remaining {gap:,} of {total_pages:,} pages ({100 * gap / total_pages:.0f}%) have NO structured "
            f"data of any tracked type at all — not even {best[0]}, the site's best-covered type. Together with the "
            f"{any_schema_pct:.0f}% above, that accounts for all {total_pages:,} crawled pages."
        )
    elif best[3] > 0 and len(insights) == 1:
        insights.append(f"{best[0]} schema is your best-covered type — present on {best[3]:,} of {total_pages:,} pages.")

    missing = [c for c in coverage if c[3] == 0 and c[2] is not None]  # curated-only, has benefit text
    missing_insights = []
    for label, field, benefit, _pages_with in missing:
        required_shapes = _SCHEMA_TYPE_REQUIRED_SHAPES.get(field)
        if required_shapes is not None:
            relevant = relevant_counts.get(field, 0)
            if relevant == 0:
                # No page on the site is even shaped for this type (e.g. no
                # careers pages at all) — not a gap, just not applicable.
                continue
            noun = _SCHEMA_TYPE_SHAPE_NOUN.get(field, label.lower())
            missing_insights.append(
                f"No {label} schema found on any of the {relevant} {noun} page(s) identified — adding it {benefit}."
            )
        else:
            # Site-wide type (Breadcrumb) — no shape to narrow to, every page qualifies.
            missing_insights.append(f"No {label} schema found on any of your {total_pages} pages — adding it {benefit}.")
        if len(missing_insights) == 2:
            break
    insights += missing_insights

    # The generic _table_slide row_cap (9, once insights are shown) silently
    # dropped whichever curated type sorted last whenever a non-curated
    # extra (e.g. Fact Check) also had real coverage, pushing the curated
    # count over 9 — confirmed live: a real Lumber report showed 9 rows
    # with "Event" (0% coverage) missing entirely, nothing on the slide
    # explaining where the rest of the picture went. The 9 curated types
    # are fixed; a couple of extra real-coverage rows is the realistic
    # case, so this leaves headroom for both while still capping well
    # short of the table pushing insights off the slide.
    return _table_slide(
        prs, "Structured Data", headers, rows, col_widths=col_widths, source="Semrush Site Audit", insights=insights,
        row_cap=min(len(coverage), 11),
    )


def add_structured_data_slide_from_crawl(prs: Presentation, schema_validation: dict):
    """Same 'Structured Data' coverage slide as add_structured_data_slide,
    but sourced from this tool's own crawl (aggregate_schema_validation)
    instead of Semrush's Site Audit export. Semrush's per-type item counts
    can be stale or wrong for a given crawl (confirmed on a live client
    site: Semrush's export showed 85% Product-schema coverage while
    Google's own Rich Results Test found zero Product items, only
    LocalBusiness + Organization) — this crawl-based path is ground truth
    since it parses the page's actual JSON-LD, so it's preferred whenever
    real crawl data is available."""
    total_pages = schema_validation.get("total_pages") or 0
    if not total_pages:
        return None

    type_coverage = schema_validation.get("type_coverage") or []
    headers = ["Schema Type", "Pages With It", "Coverage"]
    col_widths = [4.0, 2.5, 2.5]
    rows = [
        (c["type"], f"{c['pages_with_it']:,} / {total_pages:,}", f"{c['coverage_pct']}%")
        for c in type_coverage
    ]

    pages_with_schema = schema_validation.get("pages_with_schema") or 0
    any_schema_pct = 100 * pages_with_schema / total_pages
    insights = [
        f"{pages_with_schema:,} of {total_pages:,} crawled URLs ({any_schema_pct:.0f}%) have structured data "
        "(schema.org JSON-LD) implemented, per this tool's own crawl.",
    ]
    missing_types = schema_validation.get("missing_types") or []
    for m in missing_types[:2]:
        insights.append(f"{m['type']} schema is entirely missing — {m['reason']}.")
    if type_coverage:
        best = type_coverage[0]
        gap = total_pages - pages_with_schema
        if gap > 0:
            insights.append(
                f"{gap:,} of {total_pages:,} pages ({100 * gap / total_pages:.0f}%) have no structured data "
                f"of any kind — not even {best['type']}, the site's best-covered type."
            )

    return _table_slide(
        prs, "Structured Data", headers, rows, col_widths=col_widths, source="Site Audit crawl (JSON-LD)",
        insights=insights,
    )


def add_schema_validation_slide(prs: Presentation, schema_validation: dict):
    """Missing REQUIRED PROPERTIES per Google's structured-data docs — not
    just type presence/absence (the Semrush-sourced Structured Data slide
    above already covers that), but "Product pages are missing 'image' on
    12 of 15 pages" style gaps. Built from this tool's own full-site Page
    Audit crawl (see aggregate_schema_validation), since Semrush's export
    only tracks per-page type presence, never property-level completeness.

    When schema_validation also carries "gsc_rich_results" (real verdicts
    from Search Console's URL Inspection API — see gsc_service.inspect_urls,
    wired in for clients with GSC connected), those rows lead the table:
    Google's own PASS/FAIL rich-result eligibility check on real pages beats
    this tool's local rule replica, which only exists because that API needs
    site ownership this tool doesn't always have. Additive only — the local
    rule-based rows still show for every type/page GSC wasn't asked about.

    Skipped entirely if no Page Audit has been run for this client yet."""
    total_pages = schema_validation.get("total_pages") or 0
    if not total_pages:
        return None

    missing_properties = schema_validation.get("missing_properties") or []
    gsc_rich_results = schema_validation.get("gsc_rich_results") or []
    missing_types = schema_validation.get("missing_types") or []
    if not missing_properties and not gsc_rich_results and not missing_types:
        return None

    headers = ["Schema Type", "Finding", "Pages Affected"]
    col_widths = [4.0, 4.0, 4.0]

    gsc_rows = []
    gsc_pass_count = 0
    gsc_fail_count = 0
    for r in gsc_rich_results:
        verdict = r.get("verdict")
        if verdict == "PASS":
            gsc_pass_count += 1
        elif verdict == "FAIL":
            gsc_fail_count += 1
        for item in r.get("detected_items") or []:
            for sub in item.get("items") or []:
                for issue in sub.get("issues") or []:
                    gsc_rows.append((
                        f"{item.get('type')} (Google-verified)",
                        issue.get("message") or "Flagged by Google's Rich Results check",
                        _truncate_cell(r.get("url") or "", col_widths[2]),
                    ))

    type_rows = [
        (m["type"], "(entire type missing)", m["reason"])
        for m in missing_types
    ]
    rule_rows = [
        (
            m["type"] if m["severity"] == "required" else f"{m['type']} (recommended)",
            f"Missing {m['field']}",
            f"{m['pages_missing']:,} of {total_pages:,}",
        )
        for m in missing_properties
    ]
    rows = gsc_rows + type_rows + rule_rows

    pages_with_schema = schema_validation.get("pages_with_schema") or 0
    insights = []
    if gsc_rich_results:
        insights.append(
            f"Cross-checked against Search Console's URL Inspection API on {len(gsc_rich_results):,} page(s) with "
            f"structured data — Google's own rich-result eligibility check: {gsc_pass_count:,} pass, "
            f"{gsc_fail_count:,} fail."
        )
    if missing_types:
        worst_type = missing_types[0]
        insights.append(f"{worst_type['type']} schema is entirely missing — {worst_type['reason']}.")
    if missing_properties:
        insights.append(
            f"{pages_with_schema:,} of {total_pages:,} crawled pages have some structured data — this table also lists "
            "REQUIRED property gaps (block rich-result eligibility) and RECOMMENDED gaps (e.g. Organization's 'sameAs' "
            "entity links, which strengthen how AI engines and Google cite the site) that GSC wasn't checked against.",
        )
        worst = missing_properties[0]
        worst_label = "missing (required)" if worst["severity"] == "required" else "missing (recommended)"
        insights.append(
            f"{worst['type']} schema is {worst_label} '{worst['field']}' on {worst['pages_missing']:,} page(s) — "
            "the largest single local gap found."
        )

    return _table_slide(
        prs, "Schema Validator", headers, rows, col_widths=col_widths,
        source="Site Audit crawl + Search Console URL Inspection" if gsc_rich_results else "Site Audit crawl",
        insights=insights,
    )


# Canned fix per page-level issue string from technical_seo_service.py's
# crawler (see _meta_issues() and run_multi_page_audit) — (fix text,
# severity). Matches the real manual-report "Tech Fixes" slide format
# (ISSUE | WHERE | FIX), confirmed from an EJTOY audit — more actionable
# than the SEO Issues slide's issue-type rollup, which says how many pages
# but not which ones or what to do about it.
# category: "technical" = crawlability/infrastructure (page reachability,
# mobile rendering, duplicate-content prevention) vs "seo" = on-page
# content/metadata (title, description, headings) — client asked for Tech
# Fixes split into these two groups instead of one mixed list.
_PAGE_ISSUE_FIXES = {
    "Page not reachable": ("Fix the broken link/redirect, or add a 301 redirect to a working page.", "error", "technical"),
    "Missing <title> tag": ("Add a unique, keyword-relevant <title> tag (50-60 characters).", "error", "seo"),
    "Missing meta description": ("Write a unique meta description (150-160 characters) summarizing the page.", "warn", "seo"),
    "No <h1> tag found": ("Add a single <h1> heading stating the page's main topic.", "warn", "seo"),
    "Multiple <h1> tags found": ("Keep only one <h1> per page — demote extra ones to <h2>/<h3>.", "warn", "seo"),
    "Missing mobile viewport meta tag": ("Add a viewport meta tag so the page renders correctly on mobile.", "warn", "technical"),
    "Missing canonical tag": ("Add a self-referencing canonical tag to prevent duplicate-content issues.", "info", "technical"),
    "Title tag longer than 60 characters": ("Shorten the title tag so it isn't truncated in search results.", "info", "seo"),
}
_ISSUE_SEVERITY_RANK = {"error": 0, "warn": 1, "info": 2}


def _tech_fixes_scored_rows(page_audit: dict, analytics: dict | None) -> list[tuple]:
    top_pages_by_path: dict[str, int] = {}
    for p in (analytics or {}).get("top_pages", {}).get("rows", []) if analytics else []:
        path = (p.get("path") or "").rstrip("/")
        top_pages_by_path[path] = int(float(p.get("page_views", 0) or 0))

    scored_rows = []
    for page in page_audit.get("pages", []):
        path = urlparse(page.get("url", "")).path or "/"
        page_views = top_pages_by_path.get(path.rstrip("/"), 0)
        for issue in page.get("issues", []):
            fix = _PAGE_ISSUE_FIXES.get(issue)
            if not fix:
                continue
            fix_text, severity, category = fix
            scored_rows.append((_ISSUE_SEVERITY_RANK[severity], -page_views, issue, path, fix_text, page_views, category))
    scored_rows.sort(key=lambda r: (r[0], r[1]))
    return scored_rows


def _tech_fixes_category_slide(prs: Presentation, title: str, scored_rows: list[tuple]):
    if not scored_rows:
        return None
    shown = scored_rows[:9]
    col_widths = [2.7, 2.3, 7.1]
    rows = [(_truncate_cell(issue, col_widths[0]), _truncate_cell(path, col_widths[1]), fix_text) for _, _, issue, path, fix_text, _, _ in shown]
    unreachable_count = sum(1 for _, _, issue, _, _, _, _ in scored_rows if issue == "Page not reachable")
    insights = [f"{len(scored_rows)} {title.split(' — ')[-1].lower()} issue(s) found across the crawled pages" + (f", {unreachable_count} unreachable." if unreachable_count else ".")]
    if len(scored_rows) > len(shown):
        insights.append(f"Showing the {len(shown)} highest-priority — see SEO Issues for the full breakdown by type.")
    traffic_matched = [r for r in scored_rows if r[5] > 0]
    if traffic_matched:
        top_traffic = max(traffic_matched, key=lambda r: r[5])
        insights.append(
            f"\"{top_traffic[3]}\" gets real traffic ({top_traffic[5]:,} pageviews in the reporting window) "
            f"and has a \"{top_traffic[2]}\" issue — fixing this one affects real visitors, not just crawl health."
        )
    return _table_slide(
        prs, title, ["Issue", "Where", "Fix"], rows,
        col_widths=col_widths, source="Site crawl", insights=insights,
    )


def add_tech_fixes_slide(prs: Presentation, page_audit: dict | None, analytics: dict | None = None) -> list:
    """Flattens page_audit's per-page issues (up to 20 crawled pages) into
    one Issue/Where/Fix row per (page, issue) pair, worst-severity first
    (severity stays the primary sort — an error is still an error regardless
    of traffic). When real GA4 pageview data is available, ties within the
    same severity are broken by traffic — a fix on a page real visitors
    hit sorts above the same-severity fix on a page nobody visits — and the
    single highest-traffic affected page gets called out as an insight.
    Split into two slides (Technical vs SEO issues, see _PAGE_ISSUE_FIXES'
    category field) per client request — was one mixed list before.
    Sourced from our own crawl, not Semrush — Semrush's per-page x
    per-issue-type matrix export (mega_export.csv) isn't parsed at all
    currently (parked deliberately), would give a richer full-site version
    of this same idea later if ever built."""
    if not page_audit:
        return []

    scored_rows = _tech_fixes_scored_rows(page_audit, analytics)
    if not scored_rows:
        return []

    technical_rows = [r for r in scored_rows if r[6] == "technical"]
    seo_rows = [r for r in scored_rows if r[6] == "seo"]
    slides = [
        _tech_fixes_category_slide(prs, "Tech Fixes — Technical Issues", technical_rows),
        _tech_fixes_category_slide(prs, "Tech Fixes — SEO Issues", seo_rows),
    ]
    return [s for s in slides if s]


def _truncate_cell(text: str, width_in: float, size_pt: float = 11, max_lines: int = 1) -> str:
    """Truncates a table cell's text (with an ellipsis) so it fits within
    max_lines at this column width/font size. Table rows in this file get a
    fixed height (_table_slide's row_cap * Inches(0.4)), unlike every other
    text renderer here — there's no y-cursor to check against a max_y and
    break early. A real page URL or joined issue list can run 60-130+
    chars; wrapped into a narrow column that's 2-3 lines PowerPoint then
    grows the row to fit, which pushes the table's actual rendered height
    past what this file computed for it — overlapping the insights/footer
    positioned below at that computed (not actual) height. Confirmed live
    on the Priority Issues and Tech Fixes slides (long blog-post URLs and
    page paths). Same 154-chars-per-inch-at-11pt calibration as _wrap_lines."""
    chars_per_line = max(10, int(width_in * (154 / size_pt)))
    limit = chars_per_line * max_lines
    if len(text) <= limit:
        return text
    return text[: max(limit - 1, 1)].rstrip() + "…"


def _wrap_lines(text: str, width_emu, size_pt: float = 12) -> int:
    """Estimate how many lines `text` wraps to in a box of this width/font
    size — advancing y by this (instead of a fixed single-line guess) is
    what keeps a wrapped second line from rendering on top of the next
    item. 154 chars-per-inch-at-size-11 matches the calibration already
    proven in _insights_strip below."""
    width_in = max(width_emu / 914400, 0.3)
    chars_per_line = max(10, int(width_in * (154 / size_pt)))
    return max(1, -(-len(text) // chars_per_line))


def add_core_problem_slide(prs: Presentation, core_problem: dict):
    """The report's single diagnostic thesis — everything else gathered
    synthesized into one root-cause statement plus a category breakdown
    (On-page SEO / Off-page SEO / Content & Keyword Strategy), matching
    the real manual-report "Core Problem" slide format confirmed from a
    SPOTONIX audit. AI-generated (core_problem_service.py) from
    already-gathered findings, not invented — only called when a real
    thesis was returned; a category with no points is dropped entirely
    rather than padded with generic advice."""
    thesis = core_problem.get("thesis")
    if not thesis:
        return None
    categories = [c for c in (core_problem.get("categories") or []) if c.get("points")]

    slide = _blank_slide(prs)
    _content_header(slide, "Core Problem")

    _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.3))
    _textbox(
        slide, Inches(0.9), Inches(1.3), Inches(11.5), Inches(0.9), thesis,
        size=16, bold=True, color=_accent(),
    )

    if not categories:
        return slide

    col_top, col_height = Inches(2.65), Inches(4.0)
    gap = Inches(0.2)
    total_width = Inches(12.1)
    col_width = Emu(int((total_width - gap * (len(categories) - 1)) / len(categories)))
    left = Inches(0.6)
    for cat in categories:
        _card(slide, left, col_top, col_width, col_height)
        y = col_top + Inches(0.2)
        _textbox(slide, left + Inches(0.25), y, col_width - Inches(0.5), Inches(0.35), cat.get("name", ""), size=14, bold=True, color=_accent())
        rule = slide.shapes.add_shape(1, left + Inches(0.25), y + Inches(0.36), col_width - Inches(0.5), Pt(1.5))
        _fill(rule, _accent())
        rule.shadow.inherit = False
        y += Inches(0.55)
        for point in cat["points"][:5]:
            text = f"•  {point}"
            lines = _wrap_lines(text, col_width - Inches(0.5), size_pt=11.5)
            line_h = Inches(0.24)
            box = slide.shapes.add_textbox(left + Inches(0.25), y, col_width - Inches(0.5), line_h * lines)
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = text
            run.font.size = Pt(11.5)
            run.font.color.rgb = TEXT_DARK
            y += line_h * lines + Inches(0.08)
        left += col_width + gap
    return slide


def _insights_strip(slide, left, top, width, insights, title="Key Insights", max_y=None):
    """2-5 bullet takeaways mechanically derived from the slide's own data —
    no free-text generation, every line traces back to a number on the same
    slide. Returns the bottom y (Emu) after the strip.

    max_y (defaults to leaving room for the footer) hard-caps how far this
    strip may draw — a long real insight sentence at a table with many rows
    (row_cap pushes `top` down) could otherwise run past the slide bottom or
    print over the footer text with no visible break, since this was the
    one text-list renderer in the file with no truncate-rather-than-overflow
    guard; every AI-bullet slide already stops before its card boundary the
    same way this now does."""
    if not insights:
        return top
    if max_y is None:
        max_y = SLIDE_H - Inches(0.5)
    _textbox(slide, left, top, width, Inches(0.24), title.upper(), size=9.5, bold=True, color=_accent())
    y = top + Inches(0.26)
    # Width-aware wrap estimate (~14 chars/inch at size 11) — a fixed
    # chars-per-line regardless of column width caused text in narrow
    # columns (e.g. the PageSpeed sidebar) to under-reserve height and
    # overlap the next bullet.
    text_width_in = max(width - Inches(0.18), Inches(0.5)) / 914400
    chars_per_line = max(20, int(text_width_in * 14))
    for item in insights[:5]:
        lines = max(1, -(-len(item) // chars_per_line))
        line_h = Inches(0.22)
        item_h = line_h * lines + Inches(0.05)
        if y + item_h > max_y:
            break
        _icon_dot(slide, left, y + Inches(0.07), Inches(0.08), _accent())
        _textbox(slide, left + Inches(0.18), y, width - Inches(0.18), line_h * lines, item, size=11)
        y += item_h
    return y


def _draw_table(slide, headers, rows, top, col_widths=None, row_cap=None, left=None, width=None, insights=None):
    """Shared table-drawing body behind _table_slide, factored out so a
    slide needing extra content above the table (e.g. a stat card) can draw
    its own header/card and still reuse this instead of duplicating the
    table + insights-strip logic. Returns the slide's own bottom y (Emu)
    after the table (and insights strip, if any)."""
    if left is None:
        left = Inches(0.6)
    if width is None:
        width = Inches(12.1)
    if row_cap is None:
        row_cap = 9 if insights else 14
    n_cols = len(headers)
    n_rows = min(len(rows), row_cap) + 1
    height = Inches(0.4) * n_rows
    gframe = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = gframe.table
    table.first_row = False  # suppress the built-in banded-header theme so our colors apply cleanly

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_ROW_BG
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.bold = True
        para.font.color.rgb = HEADER_ROW_TEXT

    for i, row in enumerate(rows[:row_cap], start=1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = ROW_ALT if i % 2 == 0 else WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = TEXT_DARK

    bottom = top + height
    if insights:
        bottom = _insights_strip(slide, left, bottom + Inches(0.15), width, insights)
    return bottom


def _table_slide(prs, title, headers, rows, col_widths=None, source=None, insights=None, row_cap=None):
    slide = _blank_slide(prs)
    _content_header(slide, title)
    if source:
        _textbox(slide, Inches(8.3), Inches(0.3), Inches(4.5), Inches(0.4), f"Source: {source}", size=11, color=TEXT_MUTED)
    _draw_table(slide, headers, rows, Inches(1.2), col_widths=col_widths, row_cap=row_cap, insights=insights)
    return slide


# Branded/Non-Branded page-path rules, as specified by the client. Substring
# checks against the raw path, case-insensitive — deliberately literal to
# what was specified (e.g. "/blog" matches "/blog-post" too) rather than
# adding stricter boundary logic not asked for. Order doesn't matter within
# a bucket; a path matching neither list is dropped from both segmented
# slides (still counted in the overall, unsegmented Top Pages slide above).
_BRANDED_PAGE_SIGNALS = [
    "/about", "/our-team", "/careers", "/jobs", "/news", "/newsroom", "/press",
    "/contact", "/support", "/login", "/sign-in", "/app", "/demo", "/quote", "/pricing",
]
_NONBRANDED_PAGE_SIGNALS = [
    "/features", "/modules", "/capabilities", "/products/", "/solutions/", "/services/",
    "/blog", "/articles", "/insights", "/podcast", "/webinar", "/videos",
    "/case-studies", "/success-stories", "/customers",
    "/resources", "/guides", "/whitepapers", "/tools", "/calculators", "/templates",
    "/glossary", "/wiki", "/dictionary",
]


def _classify_page_branded(path: str) -> str | None:
    if (path or "").strip("/") == "":
        return "branded"  # homepage
    text = (path or "").lower()
    if any(sig in text for sig in _BRANDED_PAGE_SIGNALS):
        return "branded"
    if any(sig in text for sig in _NONBRANDED_PAGE_SIGNALS):
        return "non-branded"
    return None


def add_top_pages_branded_split_slide(
    prs: Presentation, branded_pct: float, branded_pages: list[dict], nonbranded_pct: float, nonbranded_pages: list[dict]
):
    """Branded and Non-Branded Top Pages on ONE slide (client asked these
    combined rather than as two separate slides) — two half-width panels
    side by side, each with its own compact stat strip + table. The %
    column is each page's share of its SEGMENT's Users (not Pageviews —
    client specifically wants Users% here), separate from the stat strip's
    share of the FULL site's users."""
    if not branded_pages and not nonbranded_pages:
        return None
    slide = _blank_slide(prs)
    _content_header(slide, "Top Pages — Branded vs Non-Branded")

    def _page_label(path: str) -> str:
        return f"{path} (Home Page)" if path.strip("/") == "" else path

    half_width_in = 5.9
    lefts = [Inches(0.6), Inches(0.6 + half_width_in + 0.3)]
    panels = [("BRANDED SHARE", branded_pct, branded_pages), ("NON-BRANDED SHARE", nonbranded_pct, nonbranded_pages)]
    bottoms = []

    for (label, share_pct, pages), left in zip(panels, lefts):
        if not pages:
            continue
        card_top, card_h = Inches(1.1), Inches(0.55)
        _card(slide, left, card_top, Inches(half_width_in), card_h)
        _textbox(slide, left + Inches(0.15), card_top + Inches(0.07), Inches(2.2), Inches(0.22), label, size=9, bold=True, color=TEXT_MUTED)
        _textbox(slide, left + Inches(0.15), card_top + Inches(0.26), Inches(1.2), Inches(0.28), f"{share_pct:.0f}%", size=17, bold=True, color=_accent())
        _textbox(slide, left + Inches(1.4), card_top + Inches(0.31), Inches(half_width_in) - Inches(1.5), Inches(0.22), "of total site users", size=8, color=TEXT_MUTED)

        segment_total_users = sum(int(float(p.get("active_users", 0) or 0)) for p in pages)
        rows = [
            (
                _truncate_cell(_page_label(p["path"]), 2.6),
                f"{int(float(p['page_views'])):,}",
                f"{int(float(p.get('active_users', 0) or 0)):,}",
                f"{(int(float(p.get('active_users', 0) or 0)) / segment_total_users * 100 if segment_total_users else 0):.1f}%",
            )
            for p in pages[:7]
        ]
        bottom = _draw_table(
            slide, ["Page", "Pageviews", "Users", "% of Users"], rows, card_top + card_h + Inches(0.15),
            col_widths=[2.6, 1.0, 1.0, 1.3], left=left, width=Inches(half_width_in), row_cap=7,
        )
        bottoms.append(bottom)

    insights = []
    if branded_pages:
        top_b = max(branded_pages, key=lambda p: float(p.get("active_users", 0) or 0))
        insights.append(f"Branded: \"{_page_label(top_b['path'])}\" leads by users.")
    if nonbranded_pages:
        top_nb = max(nonbranded_pages, key=lambda p: float(p.get("active_users", 0) or 0))
        insights.append(f"Non-Branded: \"{_page_label(top_nb['path'])}\" leads by users.")
    if insights and bottoms:
        _insights_strip(slide, Inches(0.6), max(bottoms) + Inches(0.25), Inches(11.9), insights)
    return slide


def add_traffic_overview_slide(prs: Presentation, analytics: dict):
    slide = _blank_slide(prs)
    _content_header(slide, "Traffic Overview")
    span = _ga4_date_span(analytics.get("date_range"))
    source_text = f"Source: Google Analytics 4 ({span})" if span else "Source: Google Analytics 4"
    _textbox(slide, Inches(8.3), Inches(0.3), Inches(4.5), Inches(0.4), source_text, size=11, color=TEXT_MUTED)

    traffic = (analytics.get("traffic_overview") or {}).get("rows", [])
    if not traffic:
        _textbox(slide, Inches(0.6), Inches(1.4), Inches(8), Inches(0.5), "No GA4 data available for this period.", size=14, color=TEXT_MUTED)
        return slide

    def total(key):
        return sum(float(r.get(key, 0) or 0) for r in traffic)

    sessions = total("sessions")
    users = total("total_users")
    pageviews = total("page_views")
    engagement_seconds = total("engagement_duration")
    avg_engagement = (sum(float(r.get("engagement_rate", 0) or 0) for r in traffic) / len(traffic)) * 100
    avg_bounce = (sum(float(r.get("bounce_rate", 0) or 0) for r in traffic) / len(traffic)) * 100

    def _fmt_duration(total_seconds: float) -> str:
        m, s = divmod(round(total_seconds), 60)
        return f"{m}m {s:02d}s"

    avg_session_duration = engagement_seconds / sessions if sessions else 0
    avg_time_on_page = engagement_seconds / pageviews if pageviews else 0

    metrics = [
        ("Sessions", f"{sessions:,.0f}"),
        ("Users", f"{users:,.0f}"),
        ("Page views", f"{pageviews:,.0f}"),
        ("Avg. engagement", f"{avg_engagement:.1f}%"),
        ("Bounce rate", f"{avg_bounce:.1f}%"),
        ("Avg. session", _fmt_duration(avg_session_duration)),
        ("Avg. time on page", _fmt_duration(avg_time_on_page)),
    ]
    # 4 cards per row instead of one fixed-width row of 5 — 7 metrics
    # (Avg. session and Avg. time on page added) no longer fit one row
    # without shrinking the value text past legibility.
    cols_per_row = 4
    gap = Inches(0.15)
    total_width = Inches(12.1)
    card_width = Emu(int((total_width - gap * (cols_per_row - 1)) / cols_per_row))
    card_height = Inches(1.35)
    row_gap = Inches(0.15)
    for i, (label, value) in enumerate(metrics):
        row, col = divmod(i, cols_per_row)
        left = Inches(0.6) + Emu(col * (card_width + gap))
        top = Inches(1.3) + Emu(row * (card_height + row_gap))
        _card(slide, left, top, card_width, card_height)
        _textbox(slide, left + Inches(0.15), top + Inches(0.15), card_width - Inches(0.3), Inches(0.4), label, size=12, color=TEXT_MUTED)
        _textbox(slide, left + Inches(0.15), top + Inches(0.55), card_width - Inches(0.3), Inches(0.7), value, size=20, bold=True, color=_accent())

    rows_used = -(-len(metrics) // cols_per_row)
    insights_top = Inches(1.3) + Emu(rows_used * (card_height + row_gap)) + Inches(0.05)
    pages_per_session = pageviews / sessions if sessions else 0
    insights = []
    if avg_bounce >= 55:
        insights.append(f"Bounce rate is {avg_bounce:.0f}% — high; visitors are leaving without exploring, check landing-page relevance and load speed.")
    elif avg_bounce <= 35:
        insights.append(f"Bounce rate is {avg_bounce:.0f}% — strong, visitors are engaging well past the landing page.")
    else:
        insights.append(f"Bounce rate is {avg_bounce:.0f}% — typical range, room to improve landing-page engagement.")
    insights.append(f"Visitors view {pages_per_session:.1f} pages per session on average.")
    if users and sessions:
        sessions_per_user = sessions / users
        if sessions_per_user > 1.3:
            insights.append(f"{sessions_per_user:.1f} sessions per user — a meaningful share of visitors are returning, not just first-time traffic.")
    _insights_strip(slide, Inches(0.6), insights_top, Inches(11.9), insights)
    return slide


def add_traffic_spike_slide(prs: Presentation, spike: dict):
    """One slide: the single biggest single-day traffic spike in the period,
    and which age/gender/country/channel segments drove it — mechanically
    detected in ga4_service.get_traffic_spike_breakdown, not a canned
    template section. Only called when a real spike was found (see that
    function's threshold)."""
    slide = _blank_slide(prs)
    _content_header(slide, "Traffic Spike Analysis")

    from datetime import date as _date

    fmt_date = _date.fromisoformat(spike["date"]).strftime("%b %d, %Y")

    _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(1.15))
    _textbox(
        slide, Inches(0.9), Inches(1.25), Inches(11.5), Inches(0.4),
        f"{fmt_date} ({spike['day_of_week']}) — {spike['sessions']:,} sessions", size=17, bold=True, color=_accent(),
    )
    _textbox(
        slide, Inches(0.9), Inches(1.65), Inches(11.5), Inches(0.4),
        f"{spike['pct_above_avg']:.0f}% above the period average of {spike['avg_sessions']:,} sessions/day.",
        size=12.5, color=TEXT_MUTED,
    )

    columns = [
        (label, rows)
        for label, rows in [
            ("Age", spike.get("by_age") or []),
            ("Gender", spike.get("by_gender") or []),
            ("Top Countries", spike.get("by_country") or []),
            ("Channel", spike.get("by_channel") or []),
        ]
        if rows
    ]
    if not columns:
        return slide

    # Width is divided evenly across however many columns actually have
    # data (2-4 in practice) rather than a fixed per-column width, so a
    # 4th column (Channel) fits without redesigning the layout, and a
    # client with only 2 populated columns still fills the row width.
    col_top, col_height = Inches(2.55), Inches(4.1)
    gap = Inches(0.2)
    total_width = Inches(12.1)
    col_width = Emu(int((total_width - gap * (len(columns) - 1)) / len(columns)))
    left = Inches(0.6)
    for label, rows in columns:
        _card(slide, left, col_top, col_width, col_height)
        y = col_top + Inches(0.2)
        _textbox(slide, left + Inches(0.25), y, col_width - Inches(0.5), Inches(0.35), label, size=14, bold=True, color=_accent())
        rule = slide.shapes.add_shape(1, left + Inches(0.25), y + Inches(0.36), col_width - Inches(0.5), Pt(1.5))
        _fill(rule, _accent())
        rule.shadow.inherit = False
        y += Inches(0.55)
        for row in rows:
            _textbox(slide, left + Inches(0.25), y, col_width - Inches(1.1), Inches(0.3), row["label"], size=12, color=TEXT_DARK)
            _textbox(
                slide, left + col_width - Inches(1.05), y, Inches(0.8), Inches(0.3),
                f"{row['pct']:.0f}%", size=12, bold=True, color=TEXT_MUTED, align=PP_ALIGN.RIGHT,
            )
            y += Inches(0.36)
        left += col_width + gap
    return slide


# Common countries abbreviated per client spec (not strict ISO — client
# asked for "US, IND, SG" specifically, a mix of 2- and 3-letter forms).
# Anything not in this table falls back to its first 3 letters uppercased,
# which is still short enough to keep the Country Split column readable.
_COUNTRY_ABBREVIATIONS = {
    "united states": "US", "united kingdom": "UK", "india": "IND", "singapore": "SG",
    "canada": "CA", "australia": "AU", "germany": "DE", "france": "FR", "spain": "ES",
    "italy": "IT", "netherlands": "NL", "ireland": "IE", "new zealand": "NZ",
    "united arab emirates": "UAE", "south africa": "ZA", "brazil": "BR", "mexico": "MX",
    "japan": "JP", "china": "CN", "south korea": "KR", "philippines": "PH",
    "indonesia": "ID", "malaysia": "MY", "vietnam": "VN", "thailand": "TH",
    "pakistan": "PK", "bangladesh": "BD", "nigeria": "NG", "kenya": "KE",
    "saudi arabia": "SA", "sweden": "SE", "norway": "NO", "denmark": "DK",
    "poland": "PL", "switzerland": "CH", "belgium": "BE", "portugal": "PT",
}


def _abbreviate_country(name: str) -> str:
    return _COUNTRY_ABBREVIATIONS.get(name.strip().lower(), name[:3].upper())


def add_traffic_channel_breakdown_slide(prs: Presentation, breakdown: dict, source: str | None = None):
    """Channel is the primary key (one row per channel, per report spec) —
    country and device are folded into that same row as each channel's own
    top-3 split, rather than getting separate sections, since channel is
    the priority lens here and country/device are secondary detail."""
    rows_data = breakdown.get("rows") or []
    if not rows_data:
        return None

    def _split_text(entries):
        return ", ".join(f"{e['label']} {e['pct']:.0f}%" for e in entries) if entries else "—"

    used_abbreviations: dict[str, str] = {}
    rows = []
    for r in rows_data:
        countries = r.get("top_countries") or []
        for c in countries:
            used_abbreviations[_abbreviate_country(c["label"])] = c["label"]
        rows.append((
            r["channel"],
            f"{r['avg_sessions_month']:,}",
            f"{r['pct_share']:.0f}%",
            _split_text([{"label": _abbreviate_country(c["label"]), "pct": c["pct"]} for c in countries]),
            _split_text([{"label": d["label"].title(), "pct": d["pct"]} for d in (r.get("top_devices") or [])]),
        ))
    months = breakdown.get("months")
    top = rows_data[0]
    insights = [
        f"{top['channel']} is the leading channel, averaging {top['avg_sessions_month']:,} sessions/month ({top['pct_share']:.0f}% of total).",
    ]
    if months:
        insights.append(f"Figures are monthly averages across the last {months:.1f} month(s) of tracked data.")
    if used_abbreviations:
        legend = ", ".join(f"{abbr} - {full}" for abbr, full in sorted(used_abbreviations.items()))
        insights.append(f"* {legend}")
    return _table_slide(
        prs, "Traffic Breakdown — Monthly Average",
        ["Channel", "Avg Sessions/mo", "% Share", "Top Countries", "Top Devices"],
        rows, col_widths=[2.0, 1.7, 1.1, 3.5, 3.8], source=source, insights=insights,
    )


def _fmt_num(v):
    """Domain Overview numeric fields come back as floats (2400.0, 616.0)
    even for whole-number counts — strip the trailing '.0' for display."""
    if v in (None, ""):
        return ""
    try:
        n = float(v)
    except (TypeError, ValueError):
        return str(v)
    return str(int(n)) if n == int(n) else f"{n:.1f}"


def add_competitor_table_slide(prs: Presentation, competitor_rows: list[dict]):
    """Matches the reference deck's Competitor Analysis table. Renders the
    full DR/Backlinks/Top Countries/Branded-split columns when the rows come
    from Domain Overview uploads; falls back to a slimmer traffic/keywords/
    common-keywords table when only an Organic Competitors export is available
    (that file type has no DR, backlinks, or branded-split data at all)."""
    has_rich_data = any(r.get("authority_score") or r.get("backlinks_total") for r in competitor_rows)
    # Domain Overview PDF's Organic Traffic/Keywords are always a single
    # country's database (confirmed: every export headed "US | Domain |
    # ..."), never Worldwide. When a domain's separate "Overview Trend" CSV
    # was uploaded with Database=Worldwide (see site_audit.py's
    # worldwide_by_domain merge), add the extra Worldwide columns — only
    # when at least one row actually has the data, so a client where nobody
    # uploaded that file yet doesn't get two permanently-blank columns.
    has_worldwide_data = any(r.get("organic_traffic_worldwide") is not None for r in competitor_rows)

    if has_rich_data and has_worldwide_data:
        headers = [
            "Domain", "Organic Traffic", "Organic Traffic (Global)", "Organic Keywords", "Organic Keywords (Global)",
            "Paid Traffic", "DR", "Backlinks", "Top Countries", "Branded", "Non-Branded",
        ]
        col_widths = [1.9, 1.0, 1.15, 1.0, 1.2, 0.95, 0.7, 1.0, 1.1, 0.9, 1.0]
        rows = [
            (
                r.get("domain", ""),
                _fmt_num(r.get("organic_traffic")),
                _fmt_num(r.get("organic_traffic_worldwide")),
                _fmt_num(r.get("organic_keywords")),
                _fmt_num(r.get("organic_keywords_worldwide")),
                _fmt_num(r.get("paid_traffic")),
                _fmt_num(r.get("authority_score")),
                _fmt_num(r.get("backlinks_total")),
                r.get("top_countries", ""),
                r.get("branded_pct", ""),
                r.get("nonbranded_pct", ""),
            )
            for r in competitor_rows[:14]
        ]
    elif has_rich_data:
        headers = ["Domain", "Organic Traffic", "Organic Keywords", "Paid Traffic", "DR", "Backlinks", "Top Countries", "Branded", "Non-Branded"]
        col_widths = [2.6, 1.3, 1.3, 1.2, 0.9, 1.2, 1.3, 1.1, 1.2]
        rows = [
            (
                r.get("domain", ""),
                _fmt_num(r.get("organic_traffic")),
                _fmt_num(r.get("organic_keywords")),
                _fmt_num(r.get("paid_traffic")),
                _fmt_num(r.get("authority_score")),
                _fmt_num(r.get("backlinks_total")),
                r.get("top_countries", ""),
                r.get("branded_pct", ""),
                r.get("nonbranded_pct", ""),
            )
            for r in competitor_rows[:14]
        ]
    else:
        headers = ["Domain", "Organic Traffic", "Organic Keywords", "Common Keywords"]
        col_widths = [5.5, 2.6, 2.6, 2.6]
        rows = [
            (
                r.get("domain", ""),
                _fmt_num(r.get("organic_traffic")),
                _fmt_num(r.get("organic_keywords")),
                _fmt_num(r.get("common_keywords")),
            )
            for r in competitor_rows[:14]
        ]

    # Own-site row is guaranteed first only for the Domain Overview path (see
    # site_audit.py's sort) — an Organic Competitors export has no "own" row
    # mixed in at all, so skip these own-vs-competitor insights there.
    own_row = competitor_rows[0] if (has_rich_data and competitor_rows) else None
    competitors = competitor_rows[1:] if own_row and len(competitor_rows) > 1 else []
    insights = []
    if own_row and competitors:
        traffic_leader = max(competitors, key=lambda r: _num(r.get("organic_traffic")))
        own_traffic = _num(own_row.get("organic_traffic"))
        leader_traffic = _num(traffic_leader.get("organic_traffic"))
        if leader_traffic > own_traffic:
            multiple = f"{leader_traffic / own_traffic:.0f}x" if own_traffic else "significantly"
            insights.append(f"{traffic_leader.get('domain', 'Top competitor')} gets {multiple} more organic traffic ({leader_traffic:,.0f} vs {own_traffic:,.0f}).")
        if has_rich_data:
            dr_leader = max(competitors, key=lambda r: _num(r.get("authority_score")))
            own_dr = _num(own_row.get("authority_score"))
            if _num(dr_leader.get("authority_score")) > own_dr:
                insights.append(f"Domain authority gap: {dr_leader.get('domain')} sits at DR {int(_num(dr_leader.get('authority_score')))} vs your {int(own_dr)}.")
            bl_leader = max(competitors, key=lambda r: _num(r.get("backlinks_total")))
            own_bl = _num(own_row.get("backlinks_total"))
            if _num(bl_leader.get("backlinks_total")) > own_bl:
                insights.append(f"{bl_leader.get('domain')} has {_num(bl_leader.get('backlinks_total')):,.0f} backlinks vs your {own_bl:,.0f}.")
        kw_leader = max(competitors, key=lambda r: _num(r.get("organic_keywords")))
        own_kw = _num(own_row.get("organic_keywords"))
        if _num(kw_leader.get("organic_keywords")) > own_kw:
            insights.append(f"{kw_leader.get('domain')} ranks for {_num(kw_leader.get('organic_keywords')) - own_kw:,.0f} more organic keywords than you.")

    own_export_date = own_row.get("export_date") if own_row else None
    own_worldwide_date_raw = own_row.get("worldwide_as_of_date") if own_row else None
    own_worldwide_date = None
    if own_worldwide_date_raw:
        from datetime import date as _date

        try:
            own_worldwide_date = _date.fromisoformat(own_worldwide_date_raw).strftime("%b %d, %Y")
        except ValueError:
            own_worldwide_date = own_worldwide_date_raw
    date_bits = [d for d in [own_export_date] if d]
    if own_worldwide_date and own_worldwide_date != own_export_date:
        date_bits.append(f"Worldwide as of {own_worldwide_date}")
    source = f"Semrush export ({'; '.join(date_bits)})" if date_bits else "Semrush export"
    return _table_slide(prs, "Competitor Analysis", headers, rows, col_widths=col_widths, source=source, insights=insights)


def add_competitor_positions_slides(prs: Presentation, competitor_positions: dict[str, list[dict]]):
    """One table slide per competitor domain from a Semrush Organic Research
    > Positions export — what that domain ranks for, and at what position.
    Matches the "Competitor Keywords: {domain}" slide format. Branded
    keywords (containing the competitor's own brand name) are excluded —
    the slide should only surface non-branded keyword opportunities."""
    slides = []
    for domain, rows in competitor_positions.items():
        brand = _brand_token(domain)
        rows = [r for r in rows if not _is_branded_keyword(r.get("keyword", ""), brand)]
        if not rows:
            continue
        sorted_rows = sorted(rows, key=lambda r: _num(r.get("search_volume")), reverse=True)
        table_rows = [
            (
                r.get("keyword", ""),
                f"{int(_num(r.get('search_volume'))):,}",
                r.get("keyword_difficulty", ""),
                r.get("position", ""),
                r.get("previous_position", ""),
            )
            for r in sorted_rows[:14]
        ]

        top1_10 = sum(1 for r in rows if 0 < _num(r.get("position")) <= 10)
        top_kw = sorted_rows[0]
        rising = [r for r in rows if 0 < _num(r.get("position")) < _num(r.get("previous_position") or r.get("position"))]
        insights = [
            f"Ranks in the top 10 for {top1_10} of {len(rows)} tracked keywords — {top1_10 / len(rows) * 100:.0f}% of their visible footprint.",
            f"Highest-volume keyword: \"{top_kw.get('keyword')}\" at position {top_kw.get('position')}, {int(_num(top_kw.get('search_volume'))):,} monthly searches.",
        ]
        if rising:
            insights.append(f"{len(rising)} keyword(s) climbing in rank — worth watching where they're pulling traffic from.")

        slide = _table_slide(
            prs, f"Competitor Keywords: {domain}",
            ["Keyword", "Search Volume", "KD", "Position", "Previous Position"], table_rows,
            col_widths=[5.0, 2.0, 1.5, 1.9, 1.9], source="Semrush — Organic Research: Positions",
            insights=insights,
        )
        slides.append(slide)
    return slides


def _num(v, default=0.0):
    try:
        return float(v)
    except (TypeError, ValueError):
        return default


def add_keyword_gap_slide(prs: Presentation, competitor_analysis: dict):
    """Table version of semrush_analysis_service's keyword-gap detection —
    the "issues" list only surfaces one summary sentence + a single top
    example; this renders the full ranked list of keywords a competitor
    ranks for that the client doesn't (or ranks far ahead on), so it reads
    like the manual report's keyword tables instead of one line of prose.
    Includes both gap types the analysis surfaces: not ranking at all, and
    ranking so far behind a page-1 competitor it's effectively invisible."""
    rows = competitor_analysis.get("keyword_gap_rows") or []
    rows = [r for r in rows if not _is_branded_keyword(r.get("keyword", ""), _brand_token(r.get("competitor_domain", "")))]
    if not rows:
        return None

    table_rows = [
        (
            r["keyword"],
            r["competitor_domain"] or "—",
            f"#{r['competitor_position']}" if r.get("competitor_position") else "—",
            f"#{r['your_position']}" if r.get("your_position") else "Not ranking",
            f"{r['search_volume']:,}",
            r["keyword_difficulty"] if r.get("keyword_difficulty") not in (None, "") else "—",
        )
        for r in rows
    ]
    total_volume = sum(r["search_volume"] for r in rows)
    top = rows[0]
    insights = [f"{len(rows)} keyword gap(s) found, {total_volume:,} combined monthly searches."]
    if top.get("competitor_domain"):
        your_pos_text = f"you're at #{top['your_position']}" if top.get("your_position") else "you don't rank at all"
        insights.append(f"Highest-volume gap: \"{top['keyword']}\" ({top['search_volume']:,} searches) — {top['competitor_domain']} ranks #{top['competitor_position']}, {your_pos_text}.")
    else:
        insights.append(f"Highest-volume gap: \"{top['keyword']}\" ({top['search_volume']:,} searches).")
    cpcs = [r["cpc"] for r in rows if r.get("cpc") not in (None, "")]
    if cpcs:
        avg_cpc = sum(cpcs) / len(cpcs)
        insights.append(f"Avg. CPC across these gaps is ${avg_cpc:.2f} — {'strong commercial intent, worth prioritizing' if avg_cpc > 10 else 'moderate commercial intent'}.")

    return _table_slide(
        prs, "Competitor Keyword Gap Analysis", ["Keyword", "Competitor", "Their Position", "Your Position", "Search Volume", "Difficulty"], table_rows,
        col_widths=[3.7, 2.4, 1.5, 1.5, 1.7, 1.1], source="Semrush Keyword Gap export", insights=insights,
    )


def add_competitor_best_at_slide(prs: Presentation, competitor_domain: str, narrative: dict):
    """Matches the reference decks' "What {Competitor} Does Well" slide —
    objective bullets on the competitor's own tactics/strengths, grounded in
    their homepage text and metrics (narrative["best_at"]), rendered ahead
    of the "Areas of Focus" response slide. Distinct from that slide: this
    one states facts about the competitor, not advice for the client.
    Absent (no slide) if the AI narrative didn't produce best_at bullets —
    same silent-skip pattern as the rest of this file's AI-derived content."""
    best_at = narrative.get("best_at") or []
    if not best_at:
        return None
    screenshot = narrative.get("screenshot")

    slide = _blank_slide(prs)
    _content_header(slide, f"What {competitor_domain} Does Well")
    card_top, card_height = Inches(1.1), Inches(5.9)
    card_width = Inches(7.8) if screenshot else Inches(12.1)
    text_width = card_width - Inches(0.75)
    _card(slide, Inches(0.6), card_top, card_width, card_height)
    y = Inches(1.35)

    if screenshot:
        img_left, img_width = Inches(8.7), Inches(3.9)
        try:
            slide.shapes.add_picture(BytesIO(screenshot), img_left, card_top, width=img_width)
        except Exception:
            pass
        else:
            _textbox(slide, img_left, card_top + Inches(2.6), img_width, Inches(0.3), competitor_domain, size=10.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    bullets_max_y = card_top + card_height - Inches(0.15)
    chars_per_line = max(20, int(text_width / 914400 * 14))
    line_h = Inches(0.24)
    for item in best_at[:6]:
        lines = max(1, -(-len(item) // chars_per_line))
        item_h = line_h * lines + Inches(0.08)
        if y + item_h > bullets_max_y:
            break
        _icon_dot(slide, Inches(0.9), y + Inches(0.08), Inches(0.09), DEFAULT_ACCENT)
        _textbox(slide, Inches(1.15), y, text_width - Inches(0.25), line_h * lines, item, size=12.5)
        y += item_h
    return slide


def _opportunity_quadrant(slide, left, top, width, height, label, items, color):
    """One quadrant card of add_competitor_opportunity_slide: a label
    header plus height-budgeted bullets, same truncate-rather-than-overflow
    discipline as the rest of this file's AI-derived content. Deliberately
    conservative chars-per-line estimate (11, vs. ~13-14 used for the wider
    single-card slides elsewhere in this file) — a narrow ~5in quadrant
    column wraps on whole words, so a flat width/avg-char-width estimate
    under-counts wrapped lines more here than it does on a wide card,
    which under-reserved height and let real (longer) wrapped text run
    into the next bullet."""
    _card(slide, left, top, width, height)
    pad = Inches(0.22)
    _textbox(slide, left + pad, top + Inches(0.15), width - pad * 2, Inches(0.3), label, size=13, bold=True, color=color)
    y = top + Inches(0.55)
    max_y = top + height - Inches(0.15)
    text_width = width - pad * 2 - Inches(0.2)
    chars_per_line = max(16, int(text_width / 914400 * 11))
    line_h = Inches(0.22)
    for item in items:
        lines = max(1, -(-len(item) // chars_per_line))
        item_h = line_h * lines + Inches(0.09)
        if y + item_h > max_y:
            break
        _icon_dot(slide, left + pad, y + Inches(0.07), Inches(0.07), color)
        _textbox(slide, left + pad + Inches(0.2), y, text_width, line_h * lines, item, size=10.5)
        y += item_h


def add_competitor_opportunity_slide(prs: Presentation, client_name: str, competitor_domain: str, narrative: dict):
    """Competitor Opportunity Analysis — the manual-deck-style bridge
    between "What {Competitor} Does Well" (evidence) and "Areas of Focus"
    (prescriptive advice): a four-quadrant WHAT COMPETITOR HAS / WHAT
    CLIENT LACKS / WHY IT MATTERS / WHAT CLIENT SHOULD BUILD layout.
    Grounded only in narrative["opportunity_analysis"] (produced in the
    same batched AI call as best_at/areas_of_focus, see
    competitor_narrative_service) — absent (no slide) if the AI didn't
    produce it, same silent-skip pattern as every other AI-derived slide
    in this file."""
    opp = narrative.get("opportunity_analysis") or {}
    has = opp.get("competitor_has") or []
    lacks = opp.get("client_lacks") or []
    why = opp.get("why_it_matters") or []
    build = opp.get("client_should_build") or []
    if not any([has, lacks, why, build]):
        return None

    slide = _blank_slide(prs)
    _content_header(slide, f"Competitor Opportunity Analysis: {competitor_domain}")

    gutter = Inches(0.3)
    left0, top0 = Inches(0.6), Inches(1.1)
    total_w, total_h = Inches(12.1), Inches(5.9)
    col_w = int((total_w - gutter) / 2)
    row_h = int((total_h - gutter) / 2)

    _opportunity_quadrant(slide, left0, top0, col_w, row_h, f"What {competitor_domain} Has", has[:4], _accent())
    _opportunity_quadrant(slide, left0 + col_w + gutter, top0, col_w, row_h, f"What {client_name} Lacks", lacks[:4], BAD)
    _opportunity_quadrant(slide, left0, top0 + row_h + gutter, col_w, row_h, "Why It Matters", why[:4], WARN)
    _opportunity_quadrant(slide, left0 + col_w + gutter, top0 + row_h + gutter, col_w, row_h, f"What {client_name} Should Build", build[:4], GOOD)
    return slide


def add_competitor_narrative_slide(prs: Presentation, client_name: str, competitor_domain: str, narrative: dict):
    """Matches the reference deck's "Areas of Focus for {Client} (vs
    {Competitor})" slide — a bulleted recommendation list followed by a
    closing "Strategic Growth Opportunity" paragraph. Always rendered in
    Cyces' own brand red (not the client's brand color) — this is
    agency-authored strategic content, not the client-branded data slides.
    Height-budgeted so long AI-generated bullets can't overflow into the
    footer: bullets stop once the budget is spent, and the closing
    paragraph is truncated with an ellipsis rather than overflowing.
    When narrative["screenshot"] (best-effort PNG bytes, may be absent if
    capture failed/was blocked) is present, narrows the text card to make
    room for the competitor's homepage screenshot on the right — matches
    the manual reference deck's visual grounding for this slide."""
    areas = narrative.get("areas_of_focus") or []
    opportunity = narrative.get("growth_opportunity")
    if not areas and not opportunity:
        return None
    screenshot = narrative.get("screenshot")

    slide = _blank_slide(prs)
    _content_header(slide, f"Areas of Focus for {client_name} (vs {competitor_domain})")
    card_top, card_height = Inches(1.1), Inches(5.9)
    card_width = Inches(7.8) if screenshot else Inches(12.1)
    text_width = card_width - Inches(0.75)
    _card(slide, Inches(0.6), card_top, card_width, card_height)
    y = Inches(1.35)

    if screenshot:
        img_left, img_width = Inches(8.7), Inches(3.9)
        try:
            slide.shapes.add_picture(BytesIO(screenshot), img_left, card_top, width=img_width)
        except Exception:
            pass  # corrupt/unreadable capture — skip the image, text side is unaffected
        else:
            _textbox(slide, img_left, card_top + Inches(2.6), img_width, Inches(0.3), competitor_domain, size=10.5, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    # Reserve room for the heading + at least 2 lines of the closing
    # paragraph before bullets are allowed to eat into that space.
    bullets_max_y = card_top + card_height - (Inches(0.7) if opportunity else Inches(0.15))
    chars_per_line = max(20, int(text_width / 914400 * 14))
    line_h = Inches(0.24)
    for item in areas[:9]:
        lines = max(1, -(-len(item) // chars_per_line))
        item_h = line_h * lines + Inches(0.08)
        if y + item_h > bullets_max_y:
            break
        _icon_dot(slide, Inches(0.9), y + Inches(0.08), Inches(0.09), DEFAULT_ACCENT)
        _textbox(slide, Inches(1.15), y, text_width - Inches(0.25), line_h * lines, item, size=12.5)
        y += item_h

    if opportunity:
        y += Inches(0.15)
        _textbox(slide, Inches(0.9), y, text_width, Inches(0.3), "Strategic Growth Opportunity:", size=13, bold=True, color=DEFAULT_ACCENT)
        y += Inches(0.34)
        chars_per_line = max(20, int(text_width / 914400 * 15))
        available_h = (card_top + card_height) - y - Inches(0.1)
        max_lines = max(1, int(available_h / Inches(0.24)))
        max_chars = max_lines * chars_per_line
        text = opportunity if len(opportunity) <= max_chars else opportunity[: max(0, max_chars - 1)].rsplit(" ", 1)[0] + "…"
        lines = max(1, -(-len(text) // chars_per_line))
        _textbox(slide, Inches(0.9), y, text_width, Inches(0.24) * lines, text, size=12, color=TEXT_DARK)
    return slide


def add_keyword_research_slide(prs: Presentation, keyword_rows: list[dict], max_clusters: int = 10):
    """One table slide per keyword cluster (Educational Toys, Development
    Skills, etc.), matching the reference deck's "Target Keywords" format —
    instead of one flat 14-row table that drowns thousands of uploaded rows
    into a single slide. Falls back to a flat table when no Cluster column
    was present in the uploaded export."""
    headers = ["Keyword", "Search Volume", "Keyword Difficulty"]
    col_widths = [7.0, 2.6, 2.5]

    clusters: dict[str, list[dict]] = {}
    for r in keyword_rows:
        label = (r.get("cluster") or "").strip()
        clusters.setdefault(label, []).append(r)

    def _keyword_insights(rows_for_group: list[dict]) -> list[str]:
        volumes = [_num(r.get("search_volume")) for r in rows_for_group]
        kds = [_num(r.get("keyword_difficulty")) for r in rows_for_group if r.get("keyword_difficulty") not in (None, "")]
        total_volume = sum(volumes)
        top = max(rows_for_group, key=lambda r: _num(r.get("search_volume")))
        easy_wins = [r for r in rows_for_group if _num(r.get("keyword_difficulty"), default=100) < 20 and _num(r.get("search_volume")) > 0]
        out = [f"{len(rows_for_group)} keywords, {total_volume:,.0f} combined monthly searches."]
        out.append(f"Top opportunity: \"{top.get('keyword')}\" — {_num(top.get('search_volume')):,.0f} searches/month, KD {top.get('keyword_difficulty', 'n/a')}.")
        if kds:
            out.append(f"Avg. keyword difficulty {sum(kds) / len(kds):.0f} — {'competitive cluster, prioritize content depth over volume' if sum(kds) / len(kds) > 40 else 'low-competition cluster, faster to rank in'}.")
        if easy_wins:
            out.append(f"{len(easy_wins)} low-difficulty (KD<20) keyword(s) with real search volume — quick-win content targets.")
        cpcs = [_num(r.get("cpc")) for r in rows_for_group if r.get("cpc") not in (None, "")]
        commercial = [r for r in rows_for_group if "commercial" in str(r.get("intent", "")).lower() or "transactional" in str(r.get("intent", "")).lower()]
        if cpcs and commercial:
            avg_cpc = sum(cpcs) / len(cpcs)
            out.append(f"Avg. CPC ${avg_cpc:.2f}, {len(commercial)} of {len(rows_for_group)} keyword(s) show commercial intent — prioritize these for conversion-focused pages.")
        elif cpcs:
            avg_cpc = sum(cpcs) / len(cpcs)
            out.append(f"Avg. CPC ${avg_cpc:.2f} — {'high commercial value, worth ranking organically for' if avg_cpc > 10 else 'moderate commercial value'}.")
        elif commercial:
            out.append(f"{len(commercial)} of {len(rows_for_group)} keyword(s) show commercial/transactional intent — prioritize these for conversion-focused pages.")
        return out

    if not clusters or set(clusters.keys()) == {""}:
        sorted_rows = sorted(keyword_rows, key=lambda r: _num(r.get("search_volume")), reverse=True)
        seen = set()
        deduped = []
        for r in sorted_rows:
            kw = r.get("keyword", "")
            if kw in seen:
                continue
            seen.add(kw)
            deduped.append(r)
        rows = [(r.get("keyword", ""), r.get("search_volume", ""), r.get("keyword_difficulty", "")) for r in deduped]
        insights = _keyword_insights(deduped) if deduped else []
        return [_table_slide(prs, "Target Keywords", headers, rows, col_widths=col_widths, source="Semrush export", insights=insights)]

    # Rank clusters by combined search volume, keep the strongest ones —
    # matches the reference deck's ~8-10 cluster slides rather than dumping
    # every long-tail cluster into its own slide.
    ranked = sorted(
        clusters.items(),
        key=lambda kv: sum(_num(r.get("search_volume")) for r in kv[1]),
        reverse=True,
    )
    slides = []
    for label, rows_for_cluster in ranked[:max_clusters]:
        sorted_rows = sorted(rows_for_cluster, key=lambda r: _num(r.get("search_volume")), reverse=True)
        seen = set()
        deduped = []
        for r in sorted_rows:
            kw = r.get("keyword", "")
            if kw in seen:
                continue
            seen.add(kw)
            deduped.append(r)
        rows = [(r.get("keyword", ""), r.get("search_volume", ""), r.get("keyword_difficulty", "")) for r in deduped]
        title = f"Target Keywords: {label}" if label else "Target Keywords"
        insights = _keyword_insights(deduped) if deduped else []
        slides.append(_table_slide(prs, title, headers, rows, col_widths=col_widths, source="Semrush export", insights=insights))
    return slides


def add_backlink_profile_slide(
    prs: Presentation, backlink_rows: list[dict], row_count: int, backlink_summary: dict | None = None,
    own_domain_rating: int | None = None,
):
    """Ahrefs/Semrush-widget-style summary — Backlinks, Referring Domains,
    Domain Rating, % dofollow, plus a Link Attributes breakdown when a
    Semrush Backlink List PDF summary was uploaded (backlink_summary). That
    export's aggregate stats are a real site-wide count, more authoritative
    than what's computed from a possibly-partial backlinks CSV, so they're
    preferred wherever both are available.

    Domain Rating itself is NOT sourced from Semrush at all (2026-08-28
    decision) — the manual report uses Ahrefs DR, a different metric on a
    different scale than Semrush's Authority Score, and Ahrefs has no free
    bulk/API access, so own_domain_rating comes from the user's manually
    entered DomainRating table, same as the Competitor Analysis DR column."""
    slide = _blank_slide(prs)
    _content_header(slide, "Backlink Profile")
    if backlink_summary and backlink_summary.get("export_date"):
        _textbox(
            slide, Inches(8.3), Inches(0.3), Inches(4.5), Inches(0.4),
            f"Source: Semrush Backlink List ({backlink_summary['export_date']})", size=11, color=TEXT_MUTED,
        )

    dofollow = sum(1 for r in backlink_rows if str(r.get("nofollow", "")).strip().lower() not in ("1", "true", "yes"))
    csv_pct_dofollow = round(100 * dofollow / len(backlink_rows)) if backlink_rows else None

    ref_domains = {urlparse(r["source_url"]).netloc for r in backlink_rows if r.get("source_url")}

    total_backlinks = backlink_summary["backlinks_total"] if backlink_summary else (row_count or None)
    total_referring_domains = backlink_summary["referring_domains"] if backlink_summary else (len(ref_domains) or None)
    pct_dofollow = backlink_summary.get("follow_pct") if backlink_summary and backlink_summary.get("follow_pct") is not None else csv_pct_dofollow
    authority_score = own_domain_rating
    authority_label = "Domain Rating"

    # Distinct from Domain Rating above (that's the site's own manually-
    # entered Ahrefs DR) — this averages each individual referring domain's
    # own authority/page score (domain_score, parsed per backlink row from
    # the Semrush CSV) to give a read on the QUALITY of sites linking in,
    # not the client's own authority.
    domain_scores = [float(r["domain_score"]) for r in backlink_rows if r.get("domain_score") not in (None, "")]
    avg_authority_score = round(sum(domain_scores) / len(domain_scores), 1) if domain_scores else None

    _card(slide, Inches(0.6), Inches(1.2), Inches(12.1), Inches(2.2))
    stats = [
        ("Backlinks", f"{int(total_backlinks):,}" if total_backlinks is not None else "—", f"{pct_dofollow:.0f}% dofollow" if pct_dofollow is not None else None),
        ("Referring domains", f"{int(total_referring_domains):,}" if total_referring_domains is not None else "—", None),
        (authority_label, str(int(authority_score)) if authority_score is not None else "—", None),
        ("Avg. authority score", f"{avg_authority_score:g}" if avg_authority_score is not None else "—", "of referring domains" if avg_authority_score is not None else None),
    ]
    card_w = Inches(2.75)
    gap = Inches(0.15)
    for i, (label, value, sub) in enumerate(stats):
        left = Inches(0.9) + Emu(i * (card_w + gap))
        _textbox(slide, left, Inches(1.4), card_w, Inches(0.35), label, size=13, color=TEXT_MUTED)
        _textbox(slide, left, Inches(1.8), card_w, Inches(0.8), value, size=30, bold=True, color=_accent())
        if sub:
            _textbox(slide, left, Inches(2.65), card_w, Inches(0.3), sub, size=12, color=TEXT_MUTED)

    y = Inches(3.65)
    attr_labels = ["Follow", "Nofollow", "Sponsored", "UGC"]
    has_attrs = backlink_summary and any(backlink_summary.get(f"{label.lower()}_count") is not None for label in attr_labels)
    if has_attrs:
        _card(slide, Inches(0.6), y, Inches(12.1), Inches(1.1))
        _textbox(slide, Inches(0.9), y + Inches(0.15), Inches(4), Inches(0.3), "Link Attributes", size=13, bold=True, color=_accent())
        attr_colors = {"follow": GOOD, "nofollow": WARN, "sponsored": TEXT_MUTED, "ugc": TEXT_MUTED}
        for i, label in enumerate(attr_labels):
            key = label.lower()
            count = backlink_summary.get(f"{key}_count")
            pct = backlink_summary.get(f"{key}_pct")
            if count is None:
                continue
            left = Inches(0.9) + Emu(i * Inches(2.9))
            _icon_dot(slide, left, y + Inches(0.62), Inches(0.11), attr_colors.get(key, TEXT_MUTED))
            _textbox(slide, left + Inches(0.22), y + Inches(0.5), Inches(2.6), Inches(0.3), f"{label}: {int(count):,} ({pct}%)", size=12, color=TEXT_DARK)
        y += Inches(1.35)
    else:
        y += Inches(0.2)

    insights = []
    if total_backlinks and total_referring_domains:
        links_per_domain = total_backlinks / total_referring_domains
        if links_per_domain > 5:
            insights.append(f"{links_per_domain:.1f} backlinks per referring domain — link profile is concentrated in a few sources, worth diversifying.")
        else:
            insights.append(f"{int(total_referring_domains)} distinct referring domains behind {int(total_backlinks)} backlinks — reasonably diverse source spread.")
    if pct_dofollow is not None:
        if pct_dofollow < 50:
            insights.append(f"Only {pct_dofollow:.0f}% of backlinks are dofollow — most links here aren't passing ranking authority.")
        else:
            insights.append(f"{pct_dofollow:.0f}% of backlinks are dofollow — the majority are passing ranking authority.")
    if authority_score is not None:
        insights.append(f"Domain Rating is {int(authority_score)} — {'a strong, established domain' if authority_score >= 40 else 'still building authority, prioritize link acquisition'}.")
    if avg_authority_score is not None:
        insights.append(f"Referring domains average an authority score of {avg_authority_score:g} — {'links are coming from generally reputable sites' if avg_authority_score >= 30 else 'link quality is on the lower end, prioritize higher-authority placements'}.")
    _insights_strip(slide, Inches(0.6), y, Inches(11.9), insights)
    return slide


# A generic, high-authority set of listing/review directories worth a
# submission for most B2B sites (matches the manual reference decks'
# "Current Brand Mentions" recommendations section — SPOTONIX itself
# recommended this same style of generic directory list, not a bespoke
# per-client set). Deliberately NOT a claim about whether the client is
# already listed anywhere — that would need a live search/citation-lookup
# API this app doesn't have, and this file's "never invent data" discipline
# rules out guessing. This is the submission-recommendation half only.
_BRAND_DIRECTORY_RECOMMENDATIONS = [
    ("G2", "Buyer-intent software marketplace — strong for B2B SaaS comparison shoppers."),
    ("Capterra", "Gartner-owned software directory — high-intent traffic, category-specific listings."),
    ("TrustRadius", "In-depth review platform enterprise buyers check before a demo call."),
    ("SoftwareSuggest", "Regional/vertical software directory — useful for reaching underserved markets."),
    ("Crunchbase", "Company profile indexed by AI/LLM training and citation sources — strengthens entity recognition."),
    ("Trustpilot", "General-purpose review platform — builds the trust signals AI Overviews and shoppers both check."),
]


def add_brand_mentions_slide(
    prs: Presentation, client_name: str, citations: list[dict] | None = None, wikipedia: dict | None = None,
):
    """Brand citation/directory-listing opportunities — the manual
    reference decks' "Current Brand Mentions" slide. When citations/
    wikipedia are supplied (see brand_citation_service — free, keyless
    DuckDuckGo/Google News + Wikipedia lookups), a "Where You're
    Already Cited" section leads the slide with grounded, real results.
    Silently falls back to directory-recommendations-only (the original
    version of this slide) when neither is available — never invents a
    citation that wasn't actually found, same discipline as the rest of
    this file."""
    citations = citations or []
    slide = _blank_slide(prs)
    _content_header(slide, "Brand Citation Opportunities")
    _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.9))
    max_y = Inches(6.9)
    y = Inches(1.3)

    if citations or wikipedia:
        _textbox(slide, Inches(0.9), y, Inches(11.5), Inches(0.3), "Where You're Already Cited", size=13, bold=True, color=_accent())
        y += Inches(0.36)
        if wikipedia:
            _icon_dot(slide, Inches(0.9), y + Inches(0.07), Inches(0.08), GOOD)
            _textbox(slide, Inches(1.15), y, Inches(11), Inches(0.3), f"Wikipedia: \"{wikipedia['title']}\"", size=11.5)
            y += Inches(0.3)
        for c in citations[:4]:
            if y > max_y - Inches(0.5):
                break
            _icon_dot(slide, Inches(0.9), y + Inches(0.07), Inches(0.08), GOOD)
            domain = urlparse(c["url"]).netloc
            _textbox(slide, Inches(1.15), y, Inches(11), Inches(0.3), f"{c['title']} — {domain}", size=11.5)
            y += Inches(0.3)
        y += Inches(0.25)

    _textbox(
        slide, Inches(0.9), y, Inches(11.5), Inches(0.5),
        "Getting listed on high-authority directories builds the trust signals both human buyers and AI answer "
        "engines check before recommending a brand.", size=12.5, color=TEXT_MUTED,
    )
    y += Inches(0.6)
    for name, blurb in _BRAND_DIRECTORY_RECOMMENDATIONS:
        if y > max_y - Inches(0.2):
            break
        _icon_dot(slide, Inches(0.9), y + Inches(0.08), Inches(0.09), _accent())
        _textbox(slide, Inches(1.15), y, Inches(2.4), Inches(0.3), name, size=13, bold=True)
        _textbox(slide, Inches(3.7), y, Inches(8.3), Inches(0.5), blurb, size=12, color=TEXT_DARK)
        y += Inches(0.62)
    _textbox(
        slide, Inches(0.9), min(y + Inches(0.1), max_y), Inches(11.3), Inches(0.4),
        f"Submit {client_name} to the directories most relevant to its category first — a targeted, complete "
        "profile beats a partial listing on every site at once.", size=11, color=TEXT_MUTED,
    )
    return slide


def _derive_next_steps(site_audit: dict | None, page_audit: dict | None) -> list[str]:
    steps = []
    if site_audit:
        for issue in site_audit.get("issues", []):
            if "https" in issue.lower():
                steps.append("Move the site fully to HTTPS — browsers flag non-HTTPS pages as insecure.")
            elif "robots" in issue.lower():
                steps.append("Add a robots.txt file so search engines can crawl the site predictably.")
            elif "sitemap" in issue.lower():
                steps.append("Publish a valid sitemap.xml and submit it in Google Search Console.")
            elif "title" in issue.lower():
                steps.append("Fix missing/oversized <title> tags — keep titles under 60 characters and unique per page.")
            elif "meta description" in issue.lower():
                steps.append("Add unique meta descriptions to every page to improve click-through from search results.")
            elif "h1" in issue.lower():
                steps.append("Ensure every page has exactly one <h1> describing its main topic.")
            elif "viewport" in issue.lower():
                steps.append("Add a mobile viewport meta tag — required for mobile usability and rankings.")
            elif "canonical" in issue.lower():
                steps.append("Add canonical tags to prevent duplicate-content issues.")
    if page_audit and page_audit.get("pages_with_issues"):
        steps.append(
            f"{page_audit['pages_with_issues']} of {page_audit['pages_checked']} crawled pages have on-page "
            "issues (missing titles/descriptions) — work through the page list and fix each."
        )
    # de-dupe while preserving order
    seen = set()
    unique = []
    for s in steps:
        if s not in seen:
            seen.add(s)
            unique.append(s)
    return unique


def add_domain_strategy_slide(prs: Presentation, domain_strategy: dict):
    """Domain Strategy finding — generic-TLD vs. single-target-country
    mismatch. Framed as a tradeoff explainer plus an open question, since
    whether to migrate depends on the client's expansion plans, which a
    crawl has no way to know."""
    slide = _blank_slide(prs)
    _content_header(slide, "Domain Strategy")
    card = _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(5.6))

    _textbox(slide, Inches(0.9), Inches(1.35), Inches(11.4), Inches(0.3), "Finding", size=13, bold=True, color=_accent())
    _textbox(slide, Inches(0.9), Inches(1.7), Inches(11.4), Inches(2.0), domain_strategy["finding"], size=12.5)

    _textbox(slide, Inches(0.9), Inches(4.0), Inches(11.4), Inches(0.3), "Open Question for the Client", size=13, bold=True, color=_accent())
    _textbox(slide, Inches(0.9), Inches(4.35), Inches(11.4), Inches(2.0), domain_strategy["open_question"], size=12.5)
    return slide


def add_ux_findings_slides(prs: Presentation, ux_findings: dict) -> list:
    """UI-Level Fixes (Issue/Where/Fix/Severity) — or, when no manual UX pass
    was done, a single slide saying so explicitly rather than silently
    skipping the dimension (report spec Rule 8). Conversion Opportunities
    (from the same ux_findings dict) now render on their own "Next Steps:
    Conversion SEO" slide instead, alongside the other Next Steps categories
    — see add_conversion_seo_next_steps_slide."""
    slides = []

    if ux_findings.get("no_ux_pass_done"):
        slide = _blank_slide(prs)
        _content_header(slide, "UI-Level Fixes")
        _card(slide, Inches(0.6), Inches(1.1), Inches(12.1), Inches(2.0))
        _textbox(slide, Inches(0.9), Inches(1.4), Inches(11.4), Inches(1.4), ux_findings["note"], size=13)
        slides.append(slide)
    elif not ux_findings.get("error"):
        fixes = ux_findings.get("ui_fixes") or []
        if fixes:
            rows = [(f.get("issue", ""), f.get("where", ""), f.get("fix", ""), f.get("severity", "")) for f in fixes]
            critical = sum(1 for f in fixes if (f.get("severity") or "").lower() == "critical")
            insights = [f"{len(fixes)} UI issue(s) found from the manual walkthrough."]
            if critical:
                insights.append(f"{critical} flagged Critical — these block a purchase or signup and should be fixed first.")
            slides.append(_table_slide(
                prs, "UI-Level Fixes", ["Issue", "Where", "Fix", "Severity"], rows,
                col_widths=[3.4, 2.8, 4.4, 1.5], source="Manual UX walkthrough", insights=insights,
            ))

    # Independent of the manual-walkthrough branch above — runs off the
    # landing page's own scraped text, so it can be present even when
    # no_ux_pass_done is True (no manual notes, but onboarding breakdown
    # still succeeded).
    onboarding = ux_findings.get("onboarding_breakdown")
    if onboarding and not onboarding.get("error"):
        onboarding_slide = add_onboarding_breakdown_slide(prs, onboarding)
        if onboarding_slide:
            slides.append(onboarding_slide)

    return slides


def add_onboarding_breakdown_slide(prs: Presentation, breakdown: dict):
    """Onboarding cognitive-bias breakdown of the landing page (Social
    Proof, Authority, Scarcity, etc.) plus the top 5 directional
    suggestions — a compact bias table on top, numbered suggestions below,
    same two-part layout style as the Backlink Profile slide's stats+list."""
    biases = breakdown.get("biases") or []
    suggestions = breakdown.get("top_suggestions") or []
    if not biases and not suggestions:
        return None

    slide = _blank_slide(prs)
    _content_header(slide, "Onboarding Breakdown")
    y = Inches(1.1)

    if biases:
        n_rows = len(biases) + 1
        row_h = Inches(0.35)
        table_h = row_h * n_rows
        left, width = Inches(0.6), Inches(12.1)
        gframe = slide.shapes.add_table(n_rows, 3, left, y, width, table_h)
        table = gframe.table
        table.first_row = False
        table.columns[0].width = Inches(2.3)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(8.6)

        for j, h in enumerate(["Bias", "Present", "Assessment"]):
            cell = table.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = HEADER_ROW_BG
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.bold = True
            para.font.color.rgb = HEADER_ROW_TEXT

        for i, b in enumerate(biases, start=1):
            present = bool(b.get("present"))
            values = [b.get("bias", ""), "Yes" if present else "No", b.get("assessment", "")]
            for j, val in enumerate(values):
                cell = table.cell(i, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT if i % 2 == 0 else WHITE
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10.5)
                para.font.color.rgb = GOOD if (j == 1 and present) else (WARN if j == 1 else TEXT_DARK)
        y += table_h + Inches(0.25)

    if suggestions:
        _textbox(slide, Inches(0.6), y, Inches(8), Inches(0.32), "Top 5 Directional Suggestions", size=14, bold=True, color=_accent())
        y += Inches(0.4)
        card_top = y
        card_bottom = Inches(6.95)
        _card(slide, Inches(0.6), card_top, Inches(12.1), card_bottom - card_top)
        y = card_top + Inches(0.18)
        for i, item in enumerate(suggestions[:5], start=1):
            lines = _wrap_lines(item, Inches(10.9), size_pt=11.5)
            line_h = Inches(0.19) * lines
            if y + line_h > card_bottom - Inches(0.1):
                break
            num = slide.shapes.add_textbox(Inches(0.85), y, Inches(0.4), Inches(0.28))
            p = num.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = f"{i}."
            r.font.bold = True
            r.font.size = Pt(11.5)
            r.font.color.rgb = _accent()
            _textbox(slide, Inches(1.25), y, Inches(10.9), line_h, item, size=11.5)
            y += line_h + Inches(0.12)
    return slide


def _next_steps_category_slide(prs: Presentation, title: str, intro: str | None, items: list[str]):
    """Shared renderer for the 4 Next Steps category slides (Local/Technical/
    Content/Conversion SEO) and the split AEO/GEO slides — one full-width
    numbered list, roomier per item than the old combined roadmap slide
    since each slide now covers only one category."""
    if not items:
        return None
    slide = _blank_slide(prs)
    _content_header(slide, title)
    y = Inches(1.05)
    if intro:
        _textbox(slide, Inches(0.6), y, Inches(12.1), Inches(0.4), intro, size=12, color=TEXT_MUTED)
        y += Inches(0.45)
    card_top = y
    card_bottom = Inches(6.7)
    _card(slide, Inches(0.6), card_top, Inches(12.1), card_bottom - card_top)
    y = card_top + Inches(0.25)
    for i, item in enumerate(items, start=1):
        lines = _wrap_lines(item, Inches(11.0), size_pt=13)
        line_h = Inches(13 * 0.02) * lines
        if y + line_h > card_bottom - Inches(0.15):
            break
        num = slide.shapes.add_textbox(Inches(0.85), y, Inches(0.4), Inches(0.3))
        p = num.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = f"{i}."
        r.font.bold = True
        r.font.size = Pt(13)
        r.font.color.rgb = _accent()
        _textbox(slide, Inches(1.25), y, Inches(11.0), line_h, item, size=13)
        y += line_h + Inches(0.22)
    return slide


def _slugify(text: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", (text or "").lower()).strip("-") or "page"


def add_programmatic_seo_slide(prs: Presentation, keyword_rows: list[dict] | None):
    """Hub + sub-page content-architecture recommendations — matches the
    manual reference deck's "Programmatic SEO Opportunities" slide (one main
    hub page per topic, sub-pages beneath it targeting specific keywords).
    Built entirely from the keyword clusters already identified for the
    Target Keywords slides — no new data source needed, just a different
    way of presenting the same clusters as a content-architecture plan."""
    if not keyword_rows:
        return None
    clusters: dict[str, list[dict]] = {}
    for r in keyword_rows:
        label = (r.get("cluster") or "").strip()
        if label:
            clusters.setdefault(label, []).append(r)
    if not clusters:
        return None

    ranked = sorted(clusters.items(), key=lambda kv: sum(_num(r.get("search_volume")) for r in kv[1]), reverse=True)
    items = []
    for label, rows_for_cluster in ranked[:6]:
        hub_slug = _slugify(label)
        top_keywords = sorted(rows_for_cluster, key=lambda r: _num(r.get("search_volume")), reverse=True)
        sub_slugs = []
        for r in top_keywords:
            slug = _slugify(r.get("keyword", ""))
            # A keyword identical to (or slugifying the same as) the cluster
            # name itself belongs on the hub page, not a redundant nested
            # duplicate of it.
            if slug and slug != hub_slug and slug not in sub_slugs:
                sub_slugs.append(slug)
            if len(sub_slugs) == 2:
                break
        if not sub_slugs:
            continue
        subpages = ", ".join(f"/{hub_slug}/{s}" for s in sub_slugs)
        items.append(f"{label}: main hub page /{hub_slug}, with sub-pages {subpages} targeting the cluster's highest-volume keywords.")
    if not items:
        return None

    intro = (
        "Content architecture built from the keyword clusters identified above — one hub page per topic, "
        "with sub-pages beneath it targeting the cluster's specific high-volume keywords."
    )
    return _next_steps_category_slide(prs, "Programmatic SEO Opportunities", intro, items)


def add_goals_slide(
    prs: Presentation,
    own_domain_rating: int | None,
    competitor_rows: list[dict] | None,
    keyword_rows: list[dict] | None,
):
    """Early Stage / Advanced Stage target slide, matching the manual
    reference deck's closing "Goal" page. Early Stage targets are derived
    from data already gathered elsewhere in the report (low-difficulty
    keywords on the table, current Domain Rating vs. the strongest tracked
    competitor's) — Advanced Stage stays qualitative/process-oriented, same
    as the reference deck's own advanced-stage bullets (SERP features, AI
    answer visibility, brand-authority signals), since those aren't
    something a crawl or export can size numerically."""
    early = []
    if keyword_rows:
        low_kd = [
            r for r in keyword_rows
            if r.get("keyword_difficulty") not in (None, "") and _num(r.get("keyword_difficulty"), default=100) < 30
        ]
        if low_kd:
            volume = sum(_num(r.get("search_volume")) for r in low_kd)
            early.append(
                f"Rank on page 1 (top 10) for the {len(low_kd)} target keyword(s) already identified under "
                f"KD 30 — {volume:,.0f} combined monthly searches on the table today."
            )
    if own_domain_rating is not None:
        competitor_drs = [
            _num(r.get("authority_score")) for r in (competitor_rows or [])
            if r.get("authority_score") not in (None, "") and r.get("domain")
        ]
        leader_dr = max(competitor_drs) if competitor_drs else None
        if leader_dr and leader_dr > own_domain_rating:
            target = min(int(leader_dr), own_domain_rating + 20)
            early.append(
                f"Increase Domain Rating from {own_domain_rating} toward {target} — the strongest tracked "
                f"competitor sits at DR {int(leader_dr)}."
            )
        else:
            early.append(f"Increase Domain Rating from {own_domain_rating} to {own_domain_rating + 15}+ through consistent, relevant backlink acquisition.")
    if not early:
        return None
    early.append("Drive consistent month-over-month organic traffic growth from the keyword and content work above.")

    advanced = [
        "Rank for high-difficulty (KD 50+) category keywords once the page-1 foundation from Early Stage is established.",
        "Diversify traffic beyond traditional search — earn visibility in AI answers (ChatGPT, Gemini, Claude) alongside classic search results.",
        "Win SERP features: featured snippets, AI Overviews, and image search placements.",
        "Build brand-authority signals — directory/citation listings, LinkedIn referral traffic, and steady backlink growth toward the category-leading Domain Rating.",
    ]
    items = [f"Early Stage: {b}" for b in early] + [f"Advanced Stage: {b}" for b in advanced]
    intro = "Near-term targets build the foundation; advanced-stage targets compound on them once page-1 rankings and a stronger Domain Rating are in place."
    return _next_steps_category_slide(prs, "SEO Goals & Targets", intro, items)


def add_technical_seo_next_steps_slide(
    prs: Presentation,
    site_audit: dict | None,
    page_audit: dict | None,
    tech_stack: dict | None,
    domain_strategy: dict | None,
):
    items = list(_derive_next_steps(site_audit, page_audit))
    if tech_stack and tech_stack.get("https") is False and not any("HTTPS" in i for i in items):
        items.append("Move the site fully to HTTPS before any further SEO work — it's a baseline ranking and trust signal.")
    if domain_strategy:
        items.insert(0, f"Decide the domain strategy first — {domain_strategy['open_question']}")
    intro = "Foundational and on-page fixes that unblock every other SEO effort — tackle these first."
    return _next_steps_category_slide(prs, "Next Steps: Technical SEO", intro, items)


def add_content_seo_next_steps_slide(prs: Presentation, keyword_rows: list[dict] | None):
    items = []
    if keyword_rows:
        # WHAT KIND of page each keyword calls for (format), separate from
        # WHAT TOPIC it belongs to (the cluster bullets below) — a client
        # can need both a landing page AND a guide for the same topic.
        category_counts: dict[str, int] = {}
        category_volume: dict[str, float] = {}
        category_examples: dict[str, list[str]] = {}
        for r in keyword_rows:
            keyword = r.get("keyword")
            if not keyword:
                continue
            category = _classify_keyword_page_category(keyword, r.get("intent"))
            if not category:
                continue
            category_counts[category] = category_counts.get(category, 0) + 1
            category_volume[category] = category_volume.get(category, 0) + _num(r.get("search_volume"))
            examples = category_examples.setdefault(category, [])
            if len(examples) < 2:
                examples.append(keyword)

        category_action = {label: action for label, _signals, action in _KEYWORD_PAGE_CATEGORIES}
        for label, count in sorted(category_counts.items(), key=lambda kv: -category_volume.get(kv[0], 0)):
            vol = category_volume.get(label, 0)
            vol_text = f", {int(vol):,} combined monthly searches" if vol else ""
            example_text = ", ".join(f"\"{e}\"" for e in category_examples.get(label, []))
            items.append(
                f"{count} target keyword(s) are {label}-shaped{vol_text} (e.g. {example_text}) — "
                f"{category_action[label]}."
            )

        cluster_counts: dict[str, int] = {}
        cluster_volume: dict[str, float] = {}
        for r in keyword_rows:
            label = (r.get("cluster") or "").strip()
            if not label:
                continue
            cluster_counts[label] = cluster_counts.get(label, 0) + 1
            cluster_volume[label] = cluster_volume.get(label, 0) + _num(r.get("search_volume"))
        for label, count in sorted(cluster_counts.items(), key=lambda kv: -cluster_volume.get(kv[0], 0)):
            vol = cluster_volume.get(label, 0)
            vol_text = f", {int(vol):,} combined monthly searches" if vol else ""
            items.append(f"Build out content for the \"{label}\" keyword cluster — {count} keyword(s) tracked{vol_text}.")
    items.append("Audit existing content for thin or outdated pages and refresh or consolidate them to strengthen topical authority.")
    items.append("Keep a content calendar built around the highest-volume clusters above so publishing stays consistent rather than one-off.")
    intro = "Where to focus content production, based on the keyword research and clustering above."
    return _next_steps_category_slide(prs, "Next Steps: Content SEO", intro, items)


def add_local_seo_next_steps_slide(prs: Presentation, structured_data_rows: list[dict] | None):
    items = []
    if structured_data_rows:
        total = len(structured_data_rows)
        with_local = sum(1 for r in structured_data_rows if _num(r.get("local_business_items")) > 0)
        if with_local == 0:
            items.append(f"No Local Business schema found on any of the {total} crawled pages — add it to enable map/business-info rich results.")
        else:
            items.append(f"Local Business schema present on {with_local} of {total} crawled pages — extend it to any location page still missing it.")
    items += [
        "Claim and fully complete the Google Business Profile — hours, categories, services, and photos all factor into local ranking.",
        "Keep Name/Address/Phone (NAP) identical across the website, Google Business Profile, and every directory listing — inconsistencies hurt local trust signals.",
        "Build a dedicated landing page per physical location or service area, each with unique local content rather than a duplicated template.",
        "Actively request and respond to Google reviews — review volume and recency are a direct local-ranking factor.",
        "Pursue local citations and backlinks from location-relevant directories, chambers of commerce, and local press.",
    ]
    intro = "Local visibility signals — mostly checklist items a crawl can't verify directly, so the schema finding above is the one grounded data point."
    return _next_steps_category_slide(prs, "Next Steps: Local SEO", intro, items)


def add_conversion_seo_next_steps_slide(
    prs: Presentation,
    ux_findings: dict | None,
    backlink_row_count: int,
):
    items = []
    if ux_findings and not ux_findings.get("error") and not ux_findings.get("no_ux_pass_done"):
        items += list(ux_findings.get("conversion_opportunities") or [])
    if backlink_row_count:
        items.append(
            f"Turn authority growth into conversions — pair the {backlink_row_count:,} tracked backlinks with clear, "
            "tested calls-to-action on the pages that authority actually lands on."
        )
    if not items:
        return None
    intro = "Turning existing traffic into leads and sales — from the manual UX walkthrough where available."
    return _next_steps_category_slide(prs, "Next Steps: Conversion SEO", intro, items)


def add_aeo_slide(prs: Presentation, site_audit: dict | None, page_audit: dict | None):
    """Answer Engine Optimization — schema/structured-data-eligibility
    recommendations for appearing in AI Overviews and answer boxes. Split
    out from the old combined AEO & GEO slide into its own slide."""
    schema_present = None
    if site_audit and site_audit.get("meta"):
        schema_present = bool(site_audit["meta"].get("structured_data_present"))
    missing_schema_pages = None
    if page_audit and page_audit.get("pages_with_issues"):
        missing_schema_pages = page_audit.get("pages_with_issues")

    items = ["Add structured FAQ sections to every key page, answering the questions customers actually ask before buying."]
    if schema_present is False:
        items.append("Homepage has no schema.org (JSON-LD) markup — add Organization, Product, and FAQ schema so AI Overviews and rich results can parse the page.")
    elif missing_schema_pages:
        items.append(f"{missing_schema_pages} crawled page(s) are missing schema markup that other pages already have — bring them in line.")
    else:
        items.append("Implement Organization, Product/Service, FAQ, and Breadcrumb schema site-wide for AI Overview eligibility.")
    items += [
        "Write concise, extractable answer blocks (2-3 sentences) near the top of key pages — this is what LLMs quote directly.",
        "Build a dedicated FAQ hub covering the full buyer journey: eligibility, pricing, process, and comparisons.",
    ]
    return _next_steps_category_slide(prs, "Answer Engine Optimization (AEO)", None, items)


def add_geo_slide(prs: Presentation):
    """Generative Engine Optimization — standard practice for AI-citation
    entity-building, since that isn't something a crawl can measure
    directly. Split out from the old combined AEO & GEO slide."""
    items = [
        "Publish authoritative content that positions the brand as a specialist in its core category — the framing AI engines reuse when answering category questions.",
        "Check whether the brand appears on the sources LLMs actually cite for this category (industry directories, comparison sites, press) — competitors already do.",
        "Create comparison-friendly content (\"X vs Y\", \"how to choose\") that AI systems can reference directly when answering evaluation queries.",
        "Interlink product pages, guides, and FAQs into topic clusters — stronger contextual relevance improves AI-driven discovery.",
    ]
    return _next_steps_category_slide(prs, "Generative Engine Optimization (GEO)", None, items)


def build_report(
    client_name: str,
    website_url: str,
    site_audit: dict | None = None,
    page_audit: dict | None = None,
    psi_mobile: dict | None = None,
    psi_desktop: dict | None = None,
    analytics: dict | None = None,
    competitor_rows: list[dict] | None = None,
    keyword_rows: list[dict] | None = None,
    backlink_rows: list[dict] | None = None,
    backlink_row_count: int = 0,
    competitor_positions: dict[str, list[dict]] | None = None,
    competitor_narratives: dict[str, dict] | None = None,
    brand_color_hex: str | None = None,
    logo_bytes: bytes | None = None,
    company_overview: dict | None = None,
    tech_stack: dict | None = None,
    competitor_analysis: dict | None = None,
    domain_strategy: dict | None = None,
    ux_findings: dict | None = None,
    site_audit_issues: list[dict] | None = None,
    structured_data_rows: list[dict] | None = None,
    site_audit_overview: dict | None = None,
    backlink_summary: dict | None = None,
    own_domain_rating: int | None = None,
    core_problem: dict | None = None,
    site_audit_pages_rows: list[dict] | None = None,
    next_steps_ai: dict | None = None,
    schema_validation: dict | None = None,
    brand_citations: list[dict] | None = None,
    brand_wikipedia: dict | None = None,
) -> bytes:
    if brand_color_hex:
        try:
            _theme["accent"] = RGBColor.from_string(brand_color_hex.lstrip("#"))
        except ValueError:
            _theme["accent"] = DEFAULT_ACCENT
    else:
        _theme["accent"] = DEFAULT_ACCENT

    domain = website_url.replace("https://", "").replace("http://", "").rstrip("/")
    _theme["footer"] = f"{client_name}  ·  {domain}"

    try:
        return _build_report(
            client_name, website_url, site_audit, page_audit, psi_mobile, psi_desktop,
            analytics, competitor_rows, keyword_rows, backlink_rows, backlink_row_count,
            company_overview, tech_stack, competitor_analysis, competitor_positions, logo_bytes,
            competitor_narratives, domain_strategy, ux_findings, site_audit_issues, site_audit_overview,
            backlink_summary, structured_data_rows, own_domain_rating, core_problem,
            site_audit_pages_rows, next_steps_ai, schema_validation,
            brand_citations, brand_wikipedia,
        )
    finally:
        _theme["footer"] = ""
        _theme["accent"] = DEFAULT_ACCENT


def _build_report(
    client_name: str,
    website_url: str,
    site_audit: dict | None,
    page_audit: dict | None,
    psi_mobile: dict | None,
    psi_desktop: dict | None,
    analytics: dict | None,
    competitor_rows: list[dict] | None,
    keyword_rows: list[dict] | None,
    backlink_rows: list[dict] | None,
    backlink_row_count: int,
    company_overview: dict | None = None,
    tech_stack: dict | None = None,
    competitor_analysis: dict | None = None,
    competitor_positions: dict[str, list[dict]] | None = None,
    logo_bytes: bytes | None = None,
    competitor_narratives: dict[str, dict] | None = None,
    domain_strategy: dict | None = None,
    ux_findings: dict | None = None,
    site_audit_issues: list[dict] | None = None,
    site_audit_overview: dict | None = None,
    backlink_summary: dict | None = None,
    structured_data_rows: list[dict] | None = None,
    own_domain_rating: int | None = None,
    core_problem: dict | None = None,
    site_audit_pages_rows: list[dict] | None = None,
    next_steps_ai: dict | None = None,
    schema_validation: dict | None = None,
    brand_citations: list[dict] | None = None,
    brand_wikipedia: dict | None = None,
) -> bytes:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs, client_name, website_url, logo_bytes=logo_bytes, analytics=analytics)

    if company_overview:
        add_company_overview_extracted_slide(prs, client_name, company_overview)
        add_solutions_products_slide(prs, company_overview)
    elif site_audit and site_audit.get("company_summary"):
        add_company_overview_slide(prs, client_name, site_audit["company_summary"])

    if domain_strategy:
        add_domain_strategy_slide(prs, domain_strategy)

    # Understanding Current Scenario section — template order: Website
    # Performance (PageSpeed) first, then the rest of the crawl-based
    # findings. Tech Stack & Hosting now renders AFTER this whole section
    # (was previously rendered before it started) — moved per the client
    # template's specified order.
    if site_audit or page_audit or psi_mobile or psi_desktop:
        add_section_slide(prs, client_name, "Understanding Current Scenario")
        if psi_mobile or psi_desktop:
            add_pagespeed_slide(prs, psi_mobile, psi_desktop)
            add_pagespeed_issues_slide(prs, psi_mobile, psi_desktop)
        if site_audit:
            add_site_health_slide(prs, site_audit, site_audit_overview)
            if site_audit_pages_rows:
                add_site_structure_slide(prs, site_audit_pages_rows)
            add_seo_issues_slide(prs, site_audit, page_audit, site_audit_issues)
            add_critical_issues_slide(prs, site_audit_issues, site_audit_pages_rows, analytics)
            add_priority_issues_slide(prs, site_audit_pages_rows, page_audit, analytics)
            add_tech_fixes_slide(prs, page_audit, analytics)
            if schema_validation and schema_validation.get("total_pages"):
                add_structured_data_slide_from_crawl(prs, schema_validation)
            elif structured_data_rows:
                add_structured_data_slide(prs, structured_data_rows, site_audit_pages_rows)
            if schema_validation:
                add_schema_validation_slide(prs, schema_validation)

    if tech_stack:
        add_tech_stack_slide(prs, tech_stack)

    if ux_findings:
        add_ux_findings_slides(prs, ux_findings)

    if backlink_rows or backlink_summary or own_domain_rating is not None:
        add_backlink_profile_slide(prs, backlink_rows or [], backlink_row_count, backlink_summary, own_domain_rating)

    add_brand_mentions_slide(prs, client_name, brand_citations, brand_wikipedia)

    if analytics:
        add_section_slide(prs, client_name, "Traffic & Search Performance")
        ga4_span = _ga4_date_span(analytics.get("date_range"))
        gsc_span = _gsc_date_span(analytics.get("date_range"))
        ga4_source = f"Google Analytics 4 ({ga4_span})" if ga4_span else "Google Analytics 4"
        gsc_source = f"Google Search Console ({gsc_span})" if gsc_span else "Google Search Console"
        if analytics.get("traffic_overview"):
            add_traffic_overview_slide(prs, analytics)
        if analytics.get("traffic_spike"):
            add_traffic_spike_slide(prs, analytics["traffic_spike"])
        if analytics.get("traffic_channel_breakdown"):
            add_traffic_channel_breakdown_slide(prs, analytics["traffic_channel_breakdown"], source=ga4_source)
        top_pages = (analytics.get("top_pages") or {}).get("rows", [])
        top_pages = [p for p in top_pages if "career" not in (p.get("path") or "").lower()]
        if top_pages:
            total_views = sum(int(float(p.get("page_views", 0) or 0)) for p in top_pages)

            def _page_label(path: str) -> str:
                # A bare "/" reads as an unlabeled empty path to anyone without
                # an SEO background — spell out that it's the home page.
                return f"{path} (Home Page)" if path.strip("/") == "" else path

            rows = [
                (
                    _truncate_cell(_page_label(p["path"]), 6.0),
                    f"{int(float(p['page_views'])):,}",
                    f"{(int(float(p['page_views'])) / total_views * 100 if total_views else 0):.1f}%",
                    f"{int(float(p.get('active_users', 0) or 0)):,}",
                )
                for p in top_pages[:14]
            ]
            top = top_pages[0]
            top_share = int(float(top["page_views"])) / total_views * 100 if total_views else 0
            insights = [
                f"\"{_page_label(top['path'])}\" is the top page, {top_share:.0f}% of all tracked pageviews.",
                f"Top {min(3, len(top_pages))} pages account for {sum(int(float(p['page_views'])) for p in top_pages[:3]) / total_views * 100:.0f}% of total traffic." if total_views else "",
            ]
            insights = [i for i in insights if i]
            _table_slide(
                prs, "Top Pages", ["Page", "Pageviews", "% of Contribution", "Users"], rows,
                col_widths=[6.0, 2.0, 2.0, 2.1], source=ga4_source, insights=insights,
            )

            # Full (unfiltered) page list for the branded/non-branded split —
            # career pages are excluded from the overall Top Pages slide
            # above as a non-signal, but the client's own branded rules
            # explicitly classify /careers as branded, so it belongs here.
            all_pages = (analytics.get("top_pages") or {}).get("rows", [])
            classified = [(p, _classify_page_branded(p.get("path") or "")) for p in all_pages]
            site_total_users = sum(int(float(p.get("active_users", 0) or 0)) for p, _ in classified)
            branded_pages = sorted(
                (p for p, c in classified if c == "branded"), key=lambda p: float(p.get("page_views", 0) or 0), reverse=True
            )
            nonbranded_pages = sorted(
                (p for p, c in classified if c == "non-branded"), key=lambda p: float(p.get("page_views", 0) or 0), reverse=True
            )
            branded_users = sum(int(float(p.get("active_users", 0) or 0)) for p in branded_pages)
            nonbranded_users = sum(int(float(p.get("active_users", 0) or 0)) for p in nonbranded_pages)
            add_top_pages_branded_split_slide(
                prs,
                branded_users / site_total_users * 100 if site_total_users else 0, branded_pages,
                nonbranded_users / site_total_users * 100 if site_total_users else 0, nonbranded_pages,
            )
        sources = (analytics.get("traffic_sources") or {}).get("rows", [])
        if sources:
            total_sessions = sum(int(float(s.get("sessions", 0) or 0)) for s in sources)
            rows = [
                (
                    s["channel"],
                    f"{int(float(s['sessions'])):,}",
                    f"{(float(s['sessions']) / total_sessions * 100 if total_sessions else 0):.1f}%",
                    f"{int(float(s.get('new_users', 0) or 0)):,}",
                    f"{int(float(s.get('returning_users', 0) or 0)):,}",
                    f"{s['return_rate_pct']:.0f}%" if s.get("return_rate_pct") is not None else "—",
                )
                for s in sources[:14]
            ]
            top_channel = max(sources, key=lambda s: float(s.get("sessions", 0) or 0))
            top_share = float(top_channel["sessions"]) / total_sessions * 100 if total_sessions else 0
            organic = next((s for s in sources if "organic search" in s["channel"].lower()), None)
            insights = [f"{top_channel['channel']} drives {top_share:.0f}% of sessions — the dominant channel by far." if top_share > 40 else f"Traffic is split fairly evenly, {top_channel['channel']} leads at {top_share:.0f}%."]
            if organic:
                organic_share = float(organic["sessions"]) / total_sessions * 100 if total_sessions else 0
                insights.append(f"Organic Search is {organic_share:.0f}% of sessions — {'a healthy share of true search-driven discovery' if organic_share > 15 else 'a thin slice, most traffic is not coming from search yet'}.")
            else:
                insights.append("No Organic Search sessions in this period — SEO isn't driving measurable traffic yet.")
            best_return_channel = max(
                (s for s in sources if s.get("return_rate_pct") is not None), key=lambda s: s["return_rate_pct"], default=None
            )
            if best_return_channel:
                insights.append(f"{best_return_channel['channel']} has the highest return rate at {best_return_channel['return_rate_pct']:.0f}% — strongest channel for repeat visitors.")
            _table_slide(
                prs, "Traffic Sources", ["Channel", "Sessions", "% of Sessions", "New Users", "Returning Users", "Return Rate"], rows,
                col_widths=[3.4, 1.8, 1.8, 1.8, 1.9, 1.4], source=ga4_source, insights=insights,
            )
        queries = (analytics.get("search_queries") or {}).get("rows", [])
        if queries:
            # The domain-derived token alone ("lumberfi" from lumberfi.com)
            # missed real branded queries built around the actual company
            # name ("lumber careers", "lumber payroll") since "lumberfi"
            # never appears in them as a whole word — confirmed on a real
            # Lumber regen, every one of those fell into Non-Branded. Match
            # on EITHER the company-name token OR the domain token so both
            # "lumber ..." and "lumberfi ..." style queries count as branded.
            brand_tokens = {t for t in (_brand_token(client_name), _brand_token(website_url)) if t}

            def _is_branded(query: str) -> bool:
                return any(_is_branded_keyword(query, token) for token in brand_tokens)

            def _query_table_slide(title: str, subset: list[dict]):
                if not subset:
                    return
                top_q = sorted(subset, key=lambda q: q.get("clicks", 0), reverse=True)[:14]
                rows = [(q["query"], q["clicks"], q["impressions"], f"{q['ctr']*100:.1f}%", f"{q['position']:.1f}") for q in top_q]
                total_clicks = sum(q.get("clicks", 0) for q in subset)
                total_impressions = sum(q.get("impressions", 0) for q in subset)
                avg_ctr = total_clicks / total_impressions * 100 if total_impressions else 0
                best_positioned = min((q for q in subset if q.get("clicks", 0) > 0), key=lambda q: q.get("position", 999), default=None)
                insights = [f"Average CTR is {avg_ctr:.1f}% across {total_impressions:,} impressions — {'strong' if avg_ctr > 3 else 'below the ~3% search-average, titles/descriptions may need work'}."]
                if best_positioned:
                    insights.append(f"Best-ranking clicked query: \"{best_positioned['query']}\" at position {best_positioned['position']:.1f}.")
                _table_slide(
                    prs, title, ["Query", "Clicks", "Impressions", "CTR", "Avg. position"], rows,
                    col_widths=[5.5, 1.5, 1.9, 1.5, 1.7], source=gsc_source, insights=insights,
                )

            branded_queries = [q for q in queries if _is_branded(q.get("query", ""))]
            nonbranded_queries = [q for q in queries if not _is_branded(q.get("query", ""))]
            _query_table_slide("Search Queries — Branded", branded_queries)
            _query_table_slide("Search Queries — Non-Branded", nonbranded_queries)

    if competitor_rows or keyword_rows or backlink_rows or backlink_summary or competitor_positions or competitor_narratives:
        add_section_slide(prs, client_name, "Competitor & Keyword Research")
        if keyword_rows:
            add_keyword_research_slide(prs, keyword_rows)
        if competitor_rows:
            add_competitor_table_slide(prs, competitor_rows)
        if competitor_positions:
            add_competitor_positions_slides(prs, competitor_positions)
        if competitor_narratives:
            for domain, narrative in competitor_narratives.items():
                if "error" not in narrative:
                    add_competitor_best_at_slide(prs, domain, narrative)
                    add_competitor_opportunity_slide(prs, client_name, domain, narrative)
                    add_competitor_narrative_slide(prs, client_name, domain, narrative)
        if competitor_analysis and competitor_analysis.get("keyword_gap_rows"):
            add_keyword_gap_slide(prs, competitor_analysis)

    if core_problem:
        add_core_problem_slide(prs, core_problem)

    add_section_slide(prs, client_name, "Next Steps")

    # AI-generated categories (next_steps_ai) take priority when present —
    # bespoke, business-aware advice grounded in this client's actual
    # products/competitors/numbers, matching how real manual-report decks
    # handle this section (confirmed against 3 references: a category that
    # doesn't fit the business, e.g. Local SEO for a national B2B SaaS with
    # no physical locations, is dropped entirely rather than padded with
    # generic checklist filler). Falls back to the static template slide
    # per-category whenever the AI call failed, or that one category came
    # back empty/malformed — never lets one bad category blank the section.
    ai_categories = (next_steps_ai or {}).get("categories") or {}
    # Kept in sync with next_steps_service.CATEGORY_TITLES by hand (a small,
    # stable map) rather than importing that module here — this file
    # otherwise has zero app.* imports, staying a pure rendering layer with
    # no network-client dependencies pulled in transitively. content_seo has
    # no entry here — that slide always renders through the deterministic
    # classifier below, never through this AI-category routing.
    _ai_category_titles = {
        "local_seo": "Next Steps: Local SEO",
        "technical_seo": "Next Steps: Technical SEO",
        "conversion_seo": "Next Steps: Conversion SEO",
        "aeo": "Answer Engine Optimization (AEO)",
        "geo": "Generative Engine Optimization (GEO)",
        "goals": "SEO Goals & Targets",
    }

    def _next_steps_slide(key: str, fallback_fn, *fallback_args):
        category = ai_categories.get(key)
        if not category:
            return fallback_fn(*fallback_args)
        if category.get("applicable") is False:
            # The AI explicitly judged this category irrelevant to this
            # business — that's a real finding, not a gap to paper over
            # with the generic fallback checklist.
            return None
        items = category.get("items")
        if not items:
            return fallback_fn(*fallback_args)
        title = _ai_category_titles.get(key, key)
        return _next_steps_category_slide(prs, title, category.get("intro") or None, items)

    _next_steps_slide("local_seo", add_local_seo_next_steps_slide, prs, structured_data_rows)
    _next_steps_slide(
        "technical_seo", add_technical_seo_next_steps_slide, prs, site_audit, page_audit, tech_stack, domain_strategy
    )
    # Always the deterministic keyword-page-category classifier below, never
    # the AI path — this needs an EXACT, guaranteed-consistent rule applied
    # every time (comparison-shaped vs. blog-shaped vs. landing-page-shaped
    # keywords), not an AI's variable phrasing of the same idea.
    add_content_seo_next_steps_slide(prs, keyword_rows)
    add_programmatic_seo_slide(prs, keyword_rows)
    _next_steps_slide("conversion_seo", add_conversion_seo_next_steps_slide, prs, ux_findings, backlink_row_count)
    _next_steps_slide("aeo", add_aeo_slide, prs, site_audit, page_audit)
    _next_steps_slide("geo", add_geo_slide, prs)
    _next_steps_slide("goals", add_goals_slide, prs, own_domain_rating, competitor_rows, keyword_rows)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
